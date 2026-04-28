"""
app.py
──────
SiteInsight · Streamlit Dashboard — Redesigned
Theme: White + Amber/Orange accents
Navigation: Left sidebar with section selector
Changes:
  1. White + amber theme throughout
  2. Left sidebar navigation, summary metrics at top
  3. Top risks show phase + activity, not just ID
  4. Detection flags as compact colour-coded table
  5. Site photos fixed (use_container_width) + 4 curated images
  6. MOM plain language + proper .docx download
  7. S-curve added to graphs tab

Run:
    python -m streamlit run app.py
"""

import os
import io
import json
import smtplib
import pandas as pd
import streamlit as st
import matplotlib.pyplot as plt
import matplotlib
matplotlib.use("Agg")

from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
from supabase import create_client, Client
from groq import Groq
from dotenv import load_dotenv
from datetime import datetime, date

load_dotenv()

SUPABASE_URL   = os.environ["SUPABASE_URL"]
SUPABASE_KEY   = os.environ["SUPABASE_KEY"]
GROQ_API_KEY   = os.environ["GROQ_API_KEY"]
EMAIL_SENDER   = os.environ.get("EMAIL_SENDER", "")
EMAIL_PASSWORD = os.environ.get("EMAIL_PASSWORD", "")
GROQ_MODEL     = "llama-3.3-70b-versatile"

# ── Theme colours ─────────────────────────────────────────────
AMBER  = "#E87722"
AMBER2 = "#F9A825"
NAVY   = "#1E2761"
RED    = "#D32F2F"
GREEN  = "#2E7D32"
GRAY   = "#666666"
LGRAY  = "#F5F5F5"

# ── Page config ───────────────────────────────────────────────
st.set_page_config(
    page_title="SiteInsight · WPR Intelligence",
    page_icon="🏗",
    layout="wide",
)

# ── Custom CSS — construction theme ──────────────────────────
st.markdown("""
<style>
  /* Construction site background with dark overlay */
  .stApp {
    background-image:
      linear-gradient(rgba(10,15,40,0.82), rgba(10,15,40,0.82)),
      url("https://images.unsplash.com/photo-1503387762-592deb58ef4e?w=1600&q=80");
    background-size: cover;
    background-attachment: fixed;
    background-position: center;
  }
  .block-container {
    padding-top: 1.5rem;
    padding-bottom: 2rem;
    background: rgba(255,255,255,0.95);
    border-radius: 12px;
    margin-top: 0.5rem;
  }

  /* Sidebar — deep navy */
  [data-testid="stSidebar"] { background-color: #0F1735; }
  [data-testid="stSidebar"] * { color: #FFFFFF !important; }
  [data-testid="stSidebar"] .stSelectbox label { color: #A8C8F0 !important; }

  /* Headers */
  h1 { color: #1E2761 !important; font-size: 1.6rem !important; font-weight: 800 !important; }
  h2 { color: #1E2761 !important; font-size: 1.25rem !important; }
  h3 { color: #E87722 !important; font-size: 1.1rem !important; font-weight: 700 !important; }

  /* Metric cards */
  [data-testid="metric-container"] {
    background: linear-gradient(135deg, #FFF8F0, #FFFDE7);
    border: 1.5px solid #E87722;
    border-radius: 10px;
    padding: 14px;
    box-shadow: 0 2px 8px rgba(232,119,34,0.12);
  }
  [data-testid="metric-container"] label {
    color: #666666 !important;
    font-size: 0.72rem !important;
    text-transform: uppercase;
    letter-spacing: 0.5px;
  }
  [data-testid="metric-container"] [data-testid="stMetricValue"] {
    color: #1E2761 !important;
    font-size: 1.7rem !important;
    font-weight: 800 !important;
  }

  /* Buttons */
  .stButton > button {
    background: linear-gradient(90deg, #E87722, #F9A825) !important;
    color: white !important;
    border: none !important;
    border-radius: 8px !important;
    font-weight: 700 !important;
    letter-spacing: 0.3px;
    box-shadow: 0 3px 10px rgba(232,119,34,0.3);
  }
  .stButton > button:hover {
    background: linear-gradient(90deg, #C86010, #E87722) !important;
    box-shadow: 0 4px 14px rgba(232,119,34,0.4);
  }

  /* Divider */
  hr { border-color: #E87722 !important; opacity: 0.35; }

  /* Alerts */
  .stAlert { border-radius: 8px !important; }

  /* Dataframe */
  [data-testid="stDataFrame"] {
    border: 1.5px solid #E87722;
    border-radius: 8px;
    box-shadow: 0 2px 6px rgba(0,0,0,0.06);
  }

  /* Expander */
  [data-testid="stExpander"] {
    border: 1px solid #E0E0E0 !important;
    border-radius: 8px !important;
    box-shadow: 0 1px 4px rgba(0,0,0,0.05);
  }

  /* Section header bar */
  .section-bar {
    background: linear-gradient(90deg, #1E2761, #E87722);
    color: white;
    padding: 10px 18px;
    border-radius: 8px;
    font-weight: 700;
    font-size: 0.9rem;
    margin-bottom: 14px;
    box-shadow: 0 2px 8px rgba(30,39,97,0.2);
  }

  /* Risk badges */
  .badge-high   { background:#FFEBEE; color:#C62828; padding:3px 10px; border-radius:5px; font-size:0.75rem; font-weight:700; border:1px solid #FFCDD2; }
  .badge-medium { background:#FFF8E1; color:#E65100; padding:3px 10px; border-radius:5px; font-size:0.75rem; font-weight:700; border:1px solid #FFE082; }
  .badge-low    { background:#E8F5E9; color:#1B5E20; padding:3px 10px; border-radius:5px; font-size:0.75rem; font-weight:700; border:1px solid #A5D6A7; }

  /* Phase section card */
  .phase-card {
    background: #F8F9FF;
    border-left: 4px solid #E87722;
    border-radius: 6px;
    padding: 12px 16px;
    margin-bottom: 10px;
  }
</style>
""", unsafe_allow_html=True)

# ── Site photos — 4 curated realistic images ─────────────────
SITE_PHOTOS = [
    {
        "url": "https://images.unsplash.com/photo-1504307651254-35680f356dfd?w=600&q=80",
        "caption": "Phase I · Basement formwork and shuttering in progress",
        "phase": "Phase I",
    },
    {
        "url": "https://images.unsplash.com/photo-1590579491624-f98f36d4c763?w=600&q=80",
        "caption": "Phase II · Reinforcement steel works and bar bending",
        "phase": "Phase II",
    },
    {
        "url": "https://images.unsplash.com/photo-1581094794329-c8112a89af12?w=600&q=80",
        "caption": "Phase III · Excavation and earthworks",
        "phase": "Phase III",
    },
    {
        "url": "https://images.unsplash.com/photo-1558618666-fcd25c85cd64?w=600&q=80",
        "caption": "Site · Waterproofing and finishing works",
        "phase": "Site",
    },
]

# ── Clients ───────────────────────────────────────────────────
@st.cache_resource
def get_supabase() -> Client:
    return create_client(SUPABASE_URL, SUPABASE_KEY)

@st.cache_resource
def get_groq():
    return Groq(api_key=GROQ_API_KEY)

supabase    = get_supabase()
groq_client = get_groq()

# ── Data loaders ──────────────────────────────────────────────
@st.cache_data(ttl=60)
def load_weeks():
    r = supabase.table("wpr_headers").select(
        "week_number,load_type").order("week_number").execute()
    return r.data or []

@st.cache_data(ttl=60)
def load_header(week):
    r = supabase.table("wpr_headers").select("*").eq("week_number", week).execute()
    return r.data[0] if r.data else {}

@st.cache_data(ttl=60)
def load_activities(week):
    r = supabase.table("wpr_activities").select("*").eq(
        "week_number", week).order("act_id").execute()
    return pd.DataFrame(r.data) if r.data else pd.DataFrame()

@st.cache_data(ttl=30)
def load_detections(week):
    r = supabase.table("detection_results").select("*").eq(
        "week_number", week).order("rule_id").execute()
    return pd.DataFrame(r.data) if r.data else pd.DataFrame()

@st.cache_data(ttl=60)
def load_narrative(week):
    r = supabase.table("ai_narratives").select("*").eq("week_number", week).execute()
    return r.data[0] if r.data else {}

@st.cache_data(ttl=60)
def load_dq_flags(week):
    r = supabase.table("dq_flags").select("*").eq("week_number", week).execute()
    return pd.DataFrame(r.data) if r.data else pd.DataFrame()

@st.cache_data(ttl=300)
def load_all_history():
    r = supabase.table("wpr_activities").select(
        "week_number,act_id,phase,activity,variance_pct,weeks_slip,"
        "cum_actual_pct,cum_planned_pct,delay_reason,is_critical_path,"
        "total_scope,responsible_person"
    ).order("week_number").execute()
    return pd.DataFrame(r.data) if r.data else pd.DataFrame()

@st.cache_data(ttl=300)
def load_baselines():
    r = supabase.table("baselines").select("*").order("id").execute()
    return pd.DataFrame(r.data) if r.data else pd.DataFrame()

# ── Groq helper ───────────────────────────────────────────────
def call_groq(prompt: str, max_tokens: int = 1500) -> str:
    response = groq_client.chat.completions.create(
        model=GROQ_MODEL,
        messages=[{"role": "user", "content": prompt}],
        temperature=0.3,
        max_tokens=max_tokens,
    )
    return response.choices[0].message.content.strip()

# ── Format narrative list fields ─────────────────────────────
def format_narrative_list(raw) -> list:
    """Convert Groq JSON list or string into clean Python list."""
    if isinstance(raw, list):
        return raw
    if isinstance(raw, str):
        raw = raw.strip()
        # Try JSON parse first
        if raw.startswith("["):
            try:
                parsed = json.loads(raw)
                if isinstance(parsed, list):
                    return [str(item).lstrip("*• ").strip() for item in parsed if item]
            except Exception:
                pass
        # Try as JSON object with text inside
        if raw.startswith("{"):
            try:
                parsed = json.loads(raw)
                # Return values from dict
                return [str(v).lstrip("*• ").strip() for v in parsed.values() if v]
            except Exception:
                pass
        # Plain text — split by newline or bullet
        lines = [l.strip().lstrip("-•*123456789.)").strip()
                 for l in raw.split("\n") if l.strip()]
        return [l for l in lines if l]
    return [str(raw)]

# ── Email helper ──────────────────────────────────────────────
def send_email(to_email: str, subject: str, body: str) -> bool:
    try:
        msg = MIMEMultipart("alternative")
        msg["Subject"] = subject
        msg["From"]    = EMAIL_SENDER
        msg["To"]      = to_email
        msg.attach(MIMEText(body, "plain"))
        with smtplib.SMTP_SSL("smtp.gmail.com", 465) as server:
            server.login(EMAIL_SENDER, EMAIL_PASSWORD)
            server.sendmail(EMAIL_SENDER, to_email, msg.as_string())
        return True
    except Exception as e:
        st.error(f"Email failed to {to_email}: {e}")
        return False

# ── Helper functions ──────────────────────────────────────────
def compute_weekly_velocity(hdf, aid):
    act_hist = hdf[hdf["act_id"] == aid].sort_values("week_number")
    if len(act_hist) < 2:
        return 0.0
    return float(act_hist["cum_actual_pct"].diff().dropna().mean())

def compute_projected_finish(current_pct, current_week, weekly_rate, total_weeks=80):
    try:
        import math
        if not weekly_rate or weekly_rate <= 0 or math.isnan(float(weekly_rate)):
            return total_weeks
        current_pct = 0.0 if (current_pct is None or math.isnan(float(current_pct))) else float(current_pct)
        remaining   = max(0.0, 1.0 - current_pct)
        result      = current_week + remaining / weekly_rate
        if math.isnan(result) or math.isinf(result):
            return total_weeks
        return int(min(result, total_weeks * 2))
    except Exception:
        return total_weeks

def compute_trend_slope(hdf, aid, lookback=4):
    act_hist = hdf[hdf["act_id"] == aid].sort_values("week_number").tail(lookback)
    if len(act_hist) < 2:
        return 0.0
    variances = act_hist["variance_pct"].values
    return float(variances[-1] - variances[0]) / len(variances)

def compute_risk(act, history_df, lookback=4):
    aid      = act["act_id"]
    slope    = compute_trend_slope(history_df, aid, lookback)
    variance = float(act.get("variance_pct", 0) or 0)
    slip     = int(act.get("weeks_slip", 0) or 0)
    critical = bool(act.get("is_critical_path", False))
    score = 0; flags = []
    if slope < -0.015:
        score += 3; flags.append(f"Worsening trend ({abs(slope)*100:.1f}%/wk)")
    if variance < -0.05:
        score += 2; flags.append(f"{abs(variance)*100:.1f}% behind plan")
    if slip >= 2:
        score += 2; flags.append(f"{slip} weeks slip")
    if critical:
        score += 2; flags.append("Critical path")
    if slope < 0 and variance < -0.03:
        score += 1; flags.append("Negative trend + delay")
    return score, flags, slope

def update_flag_status(flag_id, new_status):
    supabase.table("detection_results").update(
        {"status": new_status}).eq("id", flag_id).execute()

# ── MOM docx builder ─────────────────────────────────────────
def build_mom_docx(mom_text: str, project_name: str,
                   meeting_date, week_num: int) -> bytes:
    try:
        from docx import Document as DocxDoc
        from docx.shared import Pt, RGBColor, Inches
        from docx.enum.text import WD_ALIGN_PARAGRAPH

        doc = DocxDoc()
        doc.core_properties.author = "SiteInsight"

        # Page margins
        for section in doc.sections:
            section.top_margin    = Inches(1)
            section.bottom_margin = Inches(1)
            section.left_margin   = Inches(1.2)
            section.right_margin  = Inches(1.2)

        # Header
        h = doc.add_paragraph()
        h.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = h.add_run("MINUTES OF MEETING")
        run.bold = True
        run.font.size = Pt(16)
        run.font.color.rgb = RGBColor(0x1E, 0x27, 0x61)

        sub = doc.add_paragraph()
        sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sr = sub.add_run(f"{project_name}  ·  Week {week_num:02d}  ·  {meeting_date}")
        sr.font.size = Pt(11)
        sr.font.color.rgb = RGBColor(0xE8, 0x77, 0x22)

        doc.add_paragraph()

        # Parse and render MOM text
        for line in mom_text.split("\n"):
            line = line.strip()
            if not line:
                doc.add_paragraph()
                continue
            if line.isupper() and len(line) > 3:
                # Section heading
                p = doc.add_paragraph()
                r = p.add_run(line)
                r.bold = True
                r.font.size = Pt(11)
                r.font.color.rgb = RGBColor(0x1E, 0x27, 0x61)
                p.paragraph_format.space_before = Pt(12)
                p.paragraph_format.space_after  = Pt(4)
                # underline with border
                p.paragraph_format.border_bottom = True
            elif line.startswith("- "):
                p = doc.add_paragraph(style="List Bullet")
                r = p.add_run(line[2:])
                r.font.size = Pt(10)
                p.paragraph_format.space_after = Pt(3)
            else:
                p = doc.add_paragraph()
                r = p.add_run(line)
                r.font.size = Pt(10)
                p.paragraph_format.space_after = Pt(3)

        # Footer
        doc.add_paragraph()
        f = doc.add_paragraph()
        f.alignment = WD_ALIGN_PARAGRAPH.CENTER
        fr = f.add_run(f"Generated by SiteInsight  ·  {datetime.today().strftime('%d %b %Y')}")
        fr.font.size = Pt(8)
        fr.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
        fr.italic = True

        buf = io.BytesIO()
        doc.save(buf)
        buf.seek(0)
        return buf.getvalue()
    except ImportError:
        return None

# ════════════════════════════════════════════════════════════════
# SIDEBAR
# ════════════════════════════════════════════════════════════════
with st.sidebar:
    st.markdown("## 🏗 SiteInsight")
    st.markdown("*WPR Intelligence Platform*")
    st.divider()

    weeks_data = load_weeks()
    if not weeks_data:
        st.error("No data loaded. Run load_history.py first.")
        st.stop()

    week_options = [w["week_number"] for w in weeks_data]
    live_weeks   = [w["week_number"] for w in weeks_data
                    if w.get("load_type") == "live"]

    selected_week = st.selectbox(
        "Select week",
        options=week_options,
        index=len(week_options) - 1,
        format_func=lambda w:
            f"Week {w:02d} {'🟢' if w in live_weeks else '📁'}",
        key="week_selector"
    )
    # Clear cache when week changes
    if "last_week" not in st.session_state or st.session_state.last_week != selected_week:
        st.cache_data.clear()
        st.session_state.last_week = selected_week

    st.divider()

    # Navigation selector
    st.markdown("**Navigate to**")
    nav = st.radio(
        "nav",
        ["📋 Intelligence Report",
         "📝 Generate MOM",
         "📊 Generate PPT",
         "📈 Generate Graphs",
         "🔄 Recovery Simulator",
         "⚠️ Early Warning"],
        label_visibility="collapsed",
    )

    st.divider()
    baselines_df = load_baselines()
    if not baselines_df.empty:
        st.markdown("**Baselines**")
        for _, row in baselines_df.iterrows():
            active = "✅ " if row["is_active"] else ""
            st.caption(
                f"{active}**{row['baseline_version']}** · "
                f"ends {row['planned_end_date']}"
            )

    st.divider()
    if st.button("🔄 Refresh"):
        st.cache_data.clear()
        st.rerun()

# ── Load data ─────────────────────────────────────────────────
header     = load_header(selected_week)
activities = load_activities(selected_week)
detections = load_detections(selected_week)
narrative  = load_narrative(selected_week)
dq_flags   = load_dq_flags(selected_week)
history_df = load_all_history()

# Risk engine
risk_list = []
if not activities.empty and not history_df.empty:
    for _, act in activities.iterrows():
        score, flags, slope = compute_risk(act, history_df)
        if score > 0:
            risk_list.append({
                "act_id":      act["act_id"],
                "phase":       act.get("phase", ""),
                "activity":    act["activity"],
                "score":       score,
                "flags":       flags,
                "slope":       slope,
                "variance":    act["variance_pct"],
                "critical":    act.get("is_critical_path", False),
                "responsible": act.get("responsible_person", "—") or "—",
            })
risk_df = pd.DataFrame(risk_list)
if not risk_df.empty:
    risk_df = risk_df.sort_values("score", ascending=False)

project_name = header.get("project_name", "Horizon Commercial Tower")
baseline_ver = header.get("baseline_version", "B1")
contractor   = header.get("contractor", "—")

# ── Calculate consistent week date from project start ─────────
PROJECT_START = date(2025, 6, 2)  # Week 1 start — pre-monsoon 2025
week_date     = PROJECT_START + pd.Timedelta(weeks=selected_week - 1)
report_date   = week_date.strftime("%d %b %Y")

# ── Page header ───────────────────────────────────────────────
st.markdown(
    f"<h1>🏗 {project_name}</h1>",
    unsafe_allow_html=True
)
st.markdown(
    f"<p style='color:#666;margin-top:-12px;'>Week {selected_week:02d} &nbsp;·&nbsp; "
    f"Baseline: <b>{baseline_ver}</b> &nbsp;·&nbsp; "
    f"Contractor: <b>{contractor}</b> &nbsp;·&nbsp; "
    f"Report date: <b>{report_date}</b></p>",
    unsafe_allow_html=True
)

if header.get("baseline_revised_this_week"):
    st.warning(
        f"⚠️ Baseline revised this week · "
        f"{header.get('baseline_revision_note','')}"
    )

# ── KPI summary — always visible at top ──────────────────────
if not activities.empty:
    total_acts      = len(activities)
    behind          = int((activities["variance_pct"] < -0.05).sum())
    critical_behind = 0
    if "is_critical_path" in activities.columns:
        critical_behind = int((
            activities[activities["is_critical_path"] == True]
            ["variance_pct"] < -0.05
        ).sum())
    avg_variance = activities["variance_pct"].mean()
    max_slip     = int(activities["weeks_slip"].max()) \
                   if "weeks_slip" in activities.columns else 0

    k1, k2, k3, k4, k5 = st.columns(5)
    k1.metric("Total activities",  total_acts)
    k2.metric("Behind plan",       behind,
              delta=f"{behind/total_acts*100:.0f}% of project",
              delta_color="inverse")
    k3.metric("Critical path risk", critical_behind,
              delta="HIGH ALERT" if critical_behind > 0 else "Clear",
              delta_color="inverse" if critical_behind > 0 else "normal")
    k4.metric("Avg variance",      f"{avg_variance*100:.1f}%",
              delta_color="inverse")
    k5.metric("Max weeks slip",    f"{max_slip} wks",
              delta_color="inverse" if max_slip > 0 else "normal")

st.divider()


# ════════════════════════════════════════════════════════════════
# SECTION: INTELLIGENCE REPORT
# ════════════════════════════════════════════════════════════════
if nav == "📋 Intelligence Report":

    # ── Top Risks ─────────────────────────────────────────────
    st.markdown("### 🚨 Top Risks This Week")
    if not risk_df.empty:
        top_risks = risk_df.head(5)
        for _, r in top_risks.iterrows():
            risk_level = "HIGH" if r["score"] >= 7 else \
                         "MEDIUM" if r["score"] >= 4 else "LOW"
            badge_class = f"badge-{risk_level.lower()}"
            location = f"{r['phase']} · {r['activity']}"
            flags_str = " · ".join(r["flags"])
            resp = r["responsible"] if r["responsible"] != "—" else "Not assigned"

            col_badge, col_body = st.columns([1, 8])
            with col_badge:
                st.markdown(
                    f'<span class="{badge_class}">{risk_level}</span>',
                    unsafe_allow_html=True
                )
            with col_body:
                st.markdown(
                    f"**{location}**  \n"
                    f"_{flags_str}_  \n"
                    f"👤 Responsible: **{resp}**"
                )
            st.divider()
    else:
        st.success("No risks identified this week.")

    # ── Executive Summary ─────────────────────────────────────
    st.markdown("### 📋 Executive Summary")
    if narrative:
        exec_sum = narrative.get("executive_summary", "") or ""
        # If executive_summary contains full JSON, parse it
        if exec_sum.strip().startswith("{"):
            try:
                parsed = json.loads(exec_sum)
                exec_sum = parsed.get("executive_summary", exec_sum)
            except Exception:
                pass
        if exec_sum:
            st.info(exec_sum)

        # Trend commentary
        if not history_df.empty:
            wk_trend = (history_df.groupby("week_number")["variance_pct"]
                        .mean().reset_index())
            last3 = wk_trend.tail(3)["variance_pct"].values
            if len(last3) >= 2:
                direction = last3[-1] - last3[0]
                if direction < -0.01:
                    st.error(
                        f"📉 Project variance is **worsening** — dropped "
                        f"{abs(direction)*100:.1f}% over last 3 weeks. "
                        f"Intervention required immediately."
                    )
                elif direction > 0.01:
                    st.success(
                        f"📈 Project variance is **recovering** — improved "
                        f"{direction*100:.1f}% over last 3 weeks."
                    )
                else:
                    st.warning(
                        "⚠️ Project variance is **stable but negative** — "
                        "no recovery visible yet."
                    )
    else:
        st.info("No narrative generated yet. Run pipeline.py for this week.")

    st.divider()

    # ── Key Risks + Recommendations ──────────────────────────
    col_r, col_rec = st.columns(2)
    with col_r:
        st.markdown("### ⚠️ Key Risks")
        if narrative:
            risks = format_narrative_list(narrative.get("key_risks", []))
            for item in risks:
                if item:
                    st.markdown(f"- {item}")
        else:
            st.caption("No data available.")

    with col_rec:
        st.markdown("### ✅ Recommendations")
        if narrative:
            recs = format_narrative_list(narrative.get("recommendations", []))
            for item in recs:
                if item:
                    st.markdown(f"- {item}")
        else:
            st.caption("No data available.")

    st.divider()

    # ── Detection Flags — compact table ──────────────────────
    st.markdown(f"### 🔍 Detection Flags ({len(detections)} this week)")

    if selected_week > 1:
        prev_det = supabase.table("detection_results").select(
            "status").eq("week_number", selected_week - 1).execute()
        if prev_det.data:
            prev_df     = pd.DataFrame(prev_det.data)
            resolved    = (prev_df["status"] == "Resolved").sum()
            in_progress = (prev_df["status"] == "In Progress").sum()
            still_open  = (prev_df["status"] == "Open").sum()
            st.caption(
                f"From Week {selected_week-1:02d}: "
                f"✓ {resolved} resolved · "
                f"⏳ {in_progress} in progress · "
                f"🔴 {still_open} still open"
            )

    if detections.empty:
        st.success("No flags raised this week.")
    else:
        rule_names = {
            "R1": "Schedule Slip", "R2": "Scope Creep",
            "R3": "Stale Issue",   "R4": "Critical Path at Risk"
        }
        sev_color = {"high": "🔴", "medium": "🟡", "low": "🟢"}
        status_options = ["Open", "In Progress", "Resolved"]

        # Build display dataframe
        rows = []
        for _, h in detections.iterrows():
            act_id = h["act_id"]
            resp   = "—"
            if not activities.empty and "responsible_person" in activities.columns:
                match = activities[activities["act_id"] == act_id]
                if not match.empty:
                    rp = match.iloc[0].get("responsible_person")
                    if rp and str(rp) not in ["nan","None",""]:
                        resp = str(rp)

            # Get phase + activity for location
            location = act_id
            if not activities.empty:
                match = activities[activities["act_id"] == act_id]
                if not match.empty:
                    row_act = match.iloc[0]
                    location = f"{row_act.get('phase','')} · {row_act.get('activity','')}"

            rows.append({
                "Severity": sev_color.get(h.get("severity","low"),"⚪"),
                "Rule":     f"{h['rule_id']} · {rule_names.get(h['rule_id'],'')}",
                "Location": location,
                "Detail":   h.get("detail",""),
                "Responsible": resp,
                "Status":   h.get("status","Open") or "Open",
                "_id":      h.get("id"),
                "_rule":    h.get("rule_id"),
                "_act":     act_id,
            })

        flags_display = pd.DataFrame(rows)

        # Show compact table
        st.dataframe(
            flags_display[["Severity","Rule","Location","Responsible","Status"]],
            use_container_width=True,
            hide_index=True,
        )

        # Status update — below the table
        with st.expander("Update flag status"):
            for row in rows:
                col_loc, col_status, col_btn = st.columns([3, 1, 1])
                with col_loc:
                    st.caption(f"{row['Severity']} {row['Location']}")
                with col_status:
                    new_s = st.selectbox(
                        "s",
                        status_options,
                        index=status_options.index(row["Status"]),
                        key=f"s_{row['_id']}_{row['_rule']}_{row['_act']}",
                        label_visibility="collapsed",
                    )
                with col_btn:
                    if new_s != row["Status"] and row["_id"]:
                        if st.button("Save", key=f"b_{row['_id']}"):
                            update_flag_status(row["_id"], new_s)
                            st.cache_data.clear()
                            st.rerun()

    st.divider()

    # ── Activity Table ────────────────────────────────────────
    st.markdown("### 📊 Activity Progress")
    if not activities.empty:
        phases    = ["All"] + sorted(
            activities["phase"].dropna().unique().tolist())
        sel_phase = st.selectbox("Filter by phase", phases)
        disp      = activities.copy()
        if sel_phase != "All":
            disp = disp[disp["phase"] == sel_phase]

        show_cols = [c for c in [
            "act_id","phase","activity","unit",
            "cum_actual_pct","cum_planned_pct","variance_pct",
            "weeks_slip","delay_reason","responsible_person",
            "is_critical_path","remarks"
        ] if c in disp.columns]
        disp = disp[show_cols].copy()
        for pct_col in ["cum_actual_pct","cum_planned_pct","variance_pct"]:
            if pct_col in disp.columns:
                disp[pct_col] = (disp[pct_col]*100).round(1).astype(str)+"%"
        if "is_critical_path" in disp.columns:
            disp["is_critical_path"] = disp["is_critical_path"].map(
                {True:"✅", False:"—"})
        disp.columns = [c.replace("_"," ").title() for c in disp.columns]
        st.dataframe(disp, use_container_width=True, hide_index=True)

    st.divider()

    # ── Site Photos ───────────────────────────────────────────
    st.markdown("### 📸 Site Photos")
    photo_tab1, photo_tab2 = st.tabs(["📷 Upload Site Photos", "🖼 Representative Photos"])

    with photo_tab1:
        st.caption("Upload actual site photos — they will display below by phase.")
        up1, up2 = st.columns(2)
        uploaded_p1 = up1.file_uploader("Phase I photo", type=["jpg","jpeg","png"], key="up_p1")
        uploaded_p2 = up2.file_uploader("Phase II photo", type=["jpg","jpeg","png"], key="up_p2")
        up3, up4 = st.columns(2)
        uploaded_p3 = up3.file_uploader("Phase III photo", type=["jpg","jpeg","png"], key="up_p3")
        uploaded_p4 = up4.file_uploader("Site/General photo", type=["jpg","jpeg","png"], key="up_p4")

        uploaded = [uploaded_p1, uploaded_p2, uploaded_p3, uploaded_p4]
        captions  = ["Phase I", "Phase II", "Phase III", "Site"]
        if any(uploaded):
            st.divider()
            ucols = st.columns(4)
            for i, (up, cap) in enumerate(zip(uploaded, captions)):
                if up:
                    ucols[i].image(up, caption=cap, use_column_width=True)

    with photo_tab2:
        st.caption(
            "Representative construction photos by phase. "
            "In production, site engineer uploads actual photos via the Upload tab."
        )
        p1, p2, p3, p4 = st.columns(4)
        cols = [p1, p2, p3, p4]
        for i, photo in enumerate(SITE_PHOTOS):
            with cols[i]:
                st.image(
                    photo["url"],
                    caption=photo["caption"],
                    use_column_width=True,
                )


# ════════════════════════════════════════════════════════════════
# SECTION: GENERATE MOM
# ════════════════════════════════════════════════════════════════
elif nav == "📝 Generate MOM":
    st.markdown("### 📝 Minutes of Meeting Generator")
    st.caption(
        "Generates a professional MOM in plain language. "
        "Download as a formatted Word document."
    )

    with st.form("mom_form"):
        mc1, mc2 = st.columns(2)
        meeting_date        = mc1.date_input("Meeting date", value=date.today())
        meeting_chairperson = mc2.text_input("Chaired by", value="Project Manager")

        attendees = st.text_area(
            "Attendees (one per line)",
            value="Project Manager\nPlanning Engineer\nSite Engineer\nContractor Representative",
            height=80,
        )

        st.markdown("**Sections to include**")
        sc1, sc2, sc3 = st.columns(3)
        inc_delays   = sc1.checkbox("Delayed activities", value=True)
        inc_critical = sc2.checkbox("Critical path risks", value=True)
        inc_stale    = sc3.checkbox("Recurring issues",   value=True)
        sc4, sc5, sc6 = st.columns(3)
        inc_actions  = sc4.checkbox("Action items",       value=True)
        inc_baseline = sc5.checkbox("Baseline status",    value=True)
        inc_next     = sc6.checkbox("Next week targets",  value=True)

        gen_mom = st.form_submit_button("Generate MOM", type="primary")

    if gen_mom:
        if activities.empty:
            st.warning("No activity data for this week.")
        else:
            with st.spinner("Generating MOM via Groq..."):
                behind_acts = activities[activities["variance_pct"] < -0.05]
                det_r3 = detections[
                    detections["rule_id"] == "R3"
                ] if not detections.empty else pd.DataFrame()

                sections = (
                    (["delayed_activities"] if inc_delays   else []) +
                    (["critical_path"]      if inc_critical else []) +
                    (["stale_issues"]       if inc_stale    else []) +
                    (["action_items"]       if inc_actions  else []) +
                    (["baseline_status"]    if inc_baseline else []) +
                    (["next_week_targets"]  if inc_next     else [])
                )

                # Build phase-wise activity context
                phases_data = {}
                if not behind_acts.empty:
                    for _, r in behind_acts.iterrows():
                        phase = r.get("phase", "General") or "General"
                        rp    = r.get("responsible_person", "—") or "—"
                        is_cp = bool(r.get("is_critical_path", False))
                        slip  = int(r.get("weeks_slip", 0) or 0)
                        entry = (
                            f"  Activity: {r['activity']}"
                            f"{'  [CRITICAL PATH]' if is_cp else ''}"
                            f"\n  Delay reason: {r.get('delay_reason','—')}"
                            f"\n  Weeks slip: {slip}"
                            f"\n  Responsible: {rp}"
                        )
                        if phase not in phases_data:
                            phases_data[phase] = []
                        phases_data[phase].append(entry)

                phase_context = ""
                for phase, entries in sorted(phases_data.items()):
                    phase_context += f"\n{phase}:\n"
                    for entry in entries:
                        phase_context += entry + "\n"

                attendees_list = ', '.join(
                    [a.strip() for a in attendees.split(chr(10)) if a.strip()])

                prompt = f"""You are a senior planning engineer writing formal Minutes of Meeting for a construction project weekly review meeting.

PROJECT: {project_name}
WEEK: {selected_week} | DATE: {meeting_date} | CHAIRED BY: {meeting_chairperson}
BASELINE: {baseline_ver} | CONTRACTOR: {contractor}

DELAYED ACTIVITIES BY PHASE (source data — use this to write the MOM):
{phase_context if phase_context else 'None'}

RECURRING ISSUES (flagged in previous weeks, still unresolved):
{det_r3['detail'].tolist() if not det_r3.empty else 'None'}

PROJECT SUMMARY: {narrative.get('executive_summary','') if narrative else ''}

Generate a professional MOM using EXACTLY this structure:

MINUTES OF MEETING
Project: {project_name}
Date: {meeting_date}
Week: {selected_week}
Chaired by: {meeting_chairperson}
Attendees: {attendees_list}

PROJECT STATUS SUMMARY
[2-3 sentences on overall project health in plain language — is the project recovering or worsening? Which phases are most critical?]

PHASE-WISE REVIEW

[For EACH phase that has delayed activities, write this block:]

[PHASE NAME in CAPS — e.g. PHASE I]
Status: [One sentence — overall phase health this week]
[For each delayed activity in this phase:]
- [Activity name]: [Plain description of delay and its impact on site operations]
{'  *** CRITICAL PATH IMPACT: [Specific downstream activity affected and weeks impact on project completion] ***' if inc_critical else ''}
  Responsible: [Name from data]
  Action Required: [Specific, direct action — e.g. Arrange 2 additional shuttering gangs, Submit recovery programme, Mobilise excavator from Phase II]
  Deadline: By next site review meeting

{'RECURRING ISSUES' if inc_stale else ''}
{'[For each recurring issue: state the problem clearly, who is accountable, and what happens to the project if it remains unresolved]' if inc_stale else ''}

{'ACTION ITEMS' if inc_actions else ''}
{'[Numbered list of all actions: 1. Action description | Responsible person | Deadline]' if inc_actions else ''}

{'NEXT WEEK TARGETS' if inc_next else ''}
{'[List the 3-5 must-achieve milestones before the next site review meeting, phase-wise]' if inc_next else ''}

STRICT WRITING RULES:
- Write in plain construction language — NO percentages, NO variance numbers, NO technical codes
- Describe delays in human terms: "Basement shuttering is significantly behind schedule due to monsoon rains and labour shortage"
- Every delayed activity MUST have: what is delayed, why, who is responsible, what specific action is needed
- Critical path activities must have *** CRITICAL PATH IMPACT *** clearly marked within their phase block
- Actions must be specific and direct — not "resolve delay" but "arrange additional shuttering gang from subcontractor by Thursday"
- MOM must read like a real experienced planning engineer wrote it — professional, direct, actionable
- ALL CAPS for main section headings only (PHASE I, RECURRING ISSUES, ACTION ITEMS etc.)
"""
                mom_text = call_groq(prompt, max_tokens=2000)

            st.divider()
            st.markdown("**Generated MOM**")
            st.text_area("", value=mom_text, height=400,
                         label_visibility="collapsed")

            # Download as .docx
            docx_bytes = build_mom_docx(
                mom_text, project_name, meeting_date, selected_week)

            dl1, dl2 = st.columns(2)
            with dl1:
                if docx_bytes:
                    st.download_button(
                        "⬇️ Download MOM as Word (.docx)",
                        docx_bytes,
                        f"MOM_{project_name.replace(' ','_')}_Wk{selected_week:02d}.docx",
                        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                    )
                else:
                    st.info("Install python-docx for Word download: pip install python-docx")
            with dl2:
                st.download_button(
                    "⬇️ Download MOM as .txt",
                    mom_text,
                    f"MOM_SiteInsight_Wk{selected_week:02d}.txt",
                    "text/plain",
                )

            # Email notifications
            st.divider()
            st.markdown("**📧 Send Email Notifications**")
            st.caption("Each responsible person receives only their own action items.")

            if not EMAIL_SENDER:
                st.warning("EMAIL_SENDER not set in .env — add it to enable notifications.")
            else:
                person_actions = {}
                if not behind_acts.empty and "responsible_person" in behind_acts.columns:
                    for _, r in behind_acts.iterrows():
                        rp = r.get("responsible_person")
                        if rp and str(rp) not in ["nan","None",""]:
                            person = str(rp)
                            if person not in person_actions:
                                person_actions[person] = []
                            person_actions[person].append(
                                f"• {r['phase']} · {r['activity']} — "
                                f"{r['delay_reason']}. "
                                f"Action required by next site review meeting."
                            )

                if not person_actions:
                    st.info("No responsible persons found. Add names to Wk11+ Excel files.")
                else:
                    email_map = {}
                    for person in person_actions:
                        ec1, ec2 = st.columns([1,2])
                        ec1.write(f"👤 **{person}**")
                        email_map[person] = ec2.text_input(
                            f"Email for {person}",
                            placeholder="name@company.com",
                            key=f"email_{person}",
                            label_visibility="collapsed",
                        )

                    if st.button("📧 Send notifications", type="primary"):
                        sent = 0
                        for person, actions in person_actions.items():
                            recipient = email_map.get(person,"").strip()
                            if not recipient or "@" not in recipient:
                                st.warning(f"No valid email for {person} — skipped")
                                continue
                            body = f"""Dear {person},

Please find below your action items from the Week {selected_week:02d} site review meeting for {project_name}.

YOUR ACTION ITEMS:
{"".join([chr(10)+a for a in actions])}

All actions are to be completed by the next site review meeting.

Meeting date: {meeting_date}
Chaired by: {meeting_chairperson}

Regards,
SiteInsight · {project_name}
"""
                            subject = (
                                f"Action Items — Week {selected_week:02d} "
                                f"Site Review · {project_name}"
                            )
                            if send_email(recipient, subject, body):
                                st.success(f"✅ Sent to {person} ({recipient})")
                                sent += 1
                        if sent > 0:
                            st.success(f"Notifications sent to {sent} person(s).")


# ════════════════════════════════════════════════════════════════
# SECTION: GENERATE PPT — Composer
# ════════════════════════════════════════════════════════════════
elif nav == "📊 Generate PPT":
    st.markdown("### 📊 Presentation Composer")
    st.caption(
        "Select exactly what goes into each slide — add custom notes, "
        "set slide order, embed graphs. Download as .pptx for further editing."
    )

    # ── Step 1: Details ───────────────────────────────────────
    st.markdown("#### Step 1 — Presentation Details")
    d1, d2 = st.columns(2)
    ppt_by       = d1.text_input("Prepared by", value="Planning Team")
    ppt_audience = d2.selectbox(
        "Audience",
        ["Project Director","Client","Internal Team","Site Review Meeting"]
    )

    st.divider()

    # ── Step 2: Content Selector ──────────────────────────────
    st.markdown("#### Step 2 — Select Content for Each Slide")
    st.caption("Tick items to include. Add custom notes below each item.")

    # Cover — always included
    st.markdown("**🎯 Cover Slide** — always included")
    st.divider()

    # Executive Summary
    inc_exec = st.checkbox("📋 Executive Summary", value=True, key="ppt_exec")
    exec_note = ""
    if inc_exec:
        default_exec = narrative.get("executive_summary","") if narrative else ""
        exec_note = st.text_area(
            "Edit or add to Executive Summary",
            value=default_exec, height=80, key="exec_note_ppt"
        )
    st.divider()

    # Top Risks
    st.markdown("**🚨 Top Risks**")
    selected_risks = []
    if not risk_df.empty:
        for _, r in risk_df.head(8).iterrows():
            risk_level = "HIGH" if r["score"]>=7 else "MEDIUM" if r["score"]>=4 else "LOW"
            risk_label = f"{r['phase']} · {r['activity']} [{risk_level}]"
            rc1, rc2 = st.columns([2, 3])
            inc = rc1.checkbox(risk_label, value=(r["score"]>=7),
                               key=f"ppt_risk_{r['act_id']}")
            note = rc2.text_input(
                "Custom note", placeholder="e.g. Client informed, recovery plan submitted",
                key=f"ppt_risk_note_{r['act_id']}", label_visibility="collapsed"
            )
            if inc:
                selected_risks.append({
                    "label": risk_label,
                    "flags": " · ".join(r["flags"]),
                    "note":  note,
                    "critical": bool(r.get("critical")),
                })
    else:
        st.caption("No risks identified this week.")
    st.divider()

    # Phase Status
    st.markdown("**📊 Phase-wise Status**")
    selected_phases_ppt = []
    if not activities.empty:
        for phase in sorted(activities["phase"].dropna().unique()):
            grp = activities[activities["phase"]==phase]
            avg_actual  = grp["cum_actual_pct"].mean() * 100
            avg_planned = grp["cum_planned_pct"].mean() * 100
            behind_n    = (grp["variance_pct"] < -0.05).sum()
            ph1, ph2 = st.columns([2, 3])
            inc = ph1.checkbox(
                f"{phase} — {avg_actual:.0f}% vs {avg_planned:.0f}% planned",
                value=True, key=f"ppt_phase_{phase}"
            )
            note = ph2.text_input(
                "Custom note", placeholder="Key concern for this phase",
                key=f"ppt_phase_note_{phase}", label_visibility="collapsed"
            )
            if inc:
                selected_phases_ppt.append({
                    "phase":    phase,
                    "actual":   avg_actual,
                    "planned":  avg_planned,
                    "behind":   int(behind_n),
                    "total":    len(grp),
                    "note":     note,
                })
    st.divider()

    # Delayed Activities
    st.markdown("**⚠️ Delayed Activities**")
    selected_delays_ppt = []
    if not activities.empty:
        behind_acts_ppt = activities[activities["variance_pct"] < -0.05]
        if behind_acts_ppt.empty:
            st.caption("No delayed activities this week.")
        else:
            for _, r in behind_acts_ppt.iterrows():
                cp_tag = " ⚡CRITICAL" if r.get("is_critical_path") else ""
                act_label = f"{r['phase']} · {r['activity']}{cp_tag}"
                da1, da2 = st.columns([2, 3])
                inc = da1.checkbox(
                    act_label,
                    value=bool(r.get("is_critical_path")),
                    key=f"ppt_delay_{r['act_id']}"
                )
                note = da2.text_input(
                    "Recovery action note",
                    placeholder="e.g. Additional gang mobilised from Phase III",
                    key=f"ppt_delay_note_{r['act_id']}",
                    label_visibility="collapsed"
                )
                if inc:
                    selected_delays_ppt.append({
                        "label":    act_label,
                        "reason":   r.get("delay_reason","—"),
                        "resp":     r.get("responsible_person","—") or "—",
                        "note":     note,
                        "critical": bool(r.get("is_critical_path")),
                    })
    st.divider()

    # Graphs to embed
    st.markdown("**📈 Graphs to Embed**")
    st.caption("Selected graphs will be placed on 1–2 slides. Download pptx to rearrange.")
    gcol1, gcol2, gcol3 = st.columns(3)
    ppt_g_scurve    = gcol1.checkbox("S-Curve",              value=True,  key="ppt_gs")
    ppt_g_variance  = gcol2.checkbox("Variance Trend",       value=True,  key="ppt_gv")
    ppt_g_phase     = gcol3.checkbox("Phase-wise Completion",value=True,  key="ppt_gp")
    gcol4, gcol5, gcol6 = st.columns(3)
    ppt_g_slip      = gcol4.checkbox("Weeks Slip",           value=True,  key="ppt_gsl")
    ppt_g_delay     = gcol5.checkbox("Delay Reasons",        value=False, key="ppt_gd")
    ppt_g_critical  = gcol6.checkbox("Critical Path",        value=False, key="ppt_gc")
    st.divider()

    # Action Items
    inc_actions_ppt = st.checkbox("✅ Action Items Slide", value=True, key="ppt_actions")
    actions_note_ppt = ""
    if inc_actions_ppt:
        actions_note_ppt = st.text_input(
            "Custom note for Action Items",
            placeholder="e.g. All actions reviewed and agreed by contractor",
            key="ppt_actions_note"
        )
    st.divider()

    # Next Week Targets
    inc_next_ppt = st.checkbox("🎯 Next Week Targets Slide", value=True, key="ppt_next")
    next_note_ppt = ""
    if inc_next_ppt:
        next_note_ppt = st.text_input(
            "Custom note for Next Week Targets",
            placeholder="e.g. Focus on Basement 1 recovery",
            key="ppt_next_note"
        )

    st.divider()

    # ── Step 3: Slide Order ───────────────────────────────────
    st.markdown("#### Step 3 — Set Slide Order")
    st.caption("Enter priority number for each section (1 = appears first)")

    order_sections = []
    if inc_exec:             order_sections.append("Executive Summary")
    if selected_risks:       order_sections.append("Top Risks")
    if selected_phases_ppt:  order_sections.append("Phase Status")
    if selected_delays_ppt:  order_sections.append("Delayed Activities")
    graph_selected = any([ppt_g_scurve,ppt_g_variance,ppt_g_phase,
                          ppt_g_slip,ppt_g_delay,ppt_g_critical])
    if graph_selected:       order_sections.append("Graphs")
    if inc_actions_ppt:      order_sections.append("Action Items")
    if inc_next_ppt:         order_sections.append("Next Week Targets")

    slide_order_map = {}
    if order_sections:
        order_cols = st.columns(min(len(order_sections), 4))
        for i, sec in enumerate(order_sections):
            slide_order_map[sec] = order_cols[i % 4].number_input(
                sec, min_value=1, max_value=20, value=i+1,
                key=f"ppt_order_{sec}"
            )
        ordered_sections = sorted(order_sections, key=lambda x: slide_order_map.get(x, 99))
    else:
        ordered_sections = []

    st.divider()

    # ── Generate Button ───────────────────────────────────────
    if st.button("🎯 Generate Presentation", type="primary"):
        with st.spinner("Building slides via Groq..."):

            # Build rich context for Groq
            risks_ctx = "\n".join([
                f"- {r['label']}: {r['flags']}"
                + (f" | Note: {r['note']}" if r['note'] else "")
                for r in selected_risks
            ]) or "None selected"

            phases_ctx = "\n".join([
                f"- {p['phase']}: {p['actual']:.0f}% actual vs {p['planned']:.0f}% planned, "
                f"{p['behind']}/{p['total']} activities behind"
                + (f" | Note: {p['note']}" if p['note'] else "")
                for p in selected_phases_ppt
            ]) or "None selected"

            delays_ctx = "\n".join([
                f"- {d['label']} | Reason: {d['reason']} | Responsible: {d['resp']}"
                + (f" | Recovery: {d['note']}" if d['note'] else "")
                for d in selected_delays_ppt
            ]) or "None selected"

            prompt = f"""You are a senior planning engineer preparing a PowerPoint presentation for {ppt_audience}.

PROJECT: {project_name} | WEEK: {selected_week} | BASELINE: {baseline_ver}
PREPARED BY: {ppt_by}

SELECTED CONTENT:
Executive Summary: {exec_note if inc_exec else 'Not included'}
Top Risks: {risks_ctx}
Phase Status: {phases_ctx}
Delayed Activities: {delays_ctx}
Action Items note: {actions_note_ppt if inc_actions_ppt else 'Not included'}
Next Week note: {next_note_ppt if inc_next_ppt else 'Not included'}

Generate slide content for these sections in order: {', '.join(ordered_sections)}

Return ONLY valid JSON — no markdown fences, no extra text:
{{
  "slides": [
    {{
      "section": "section name from ordered list",
      "title": "Slide title — specific and impactful",
      "bullets": ["bullet 1 — max 18 words, specific", "bullet 2", "bullet 3", "bullet 4", "bullet 5"],
      "highlight": "Key stat or number for this slide"
    }}
  ]
}}

Rules:
- 4-6 bullets per slide — substantive, not generic
- Use actual project data — phase names, activity names, responsible persons from context
- For {ppt_audience}: focus on {"decisions needed, milestone impact, cost risk" if ppt_audience=="Project Director" else "progress confidence, deliverables, timeline" if ppt_audience=="Client" else "who does what, specific actions, phase-wise accountability" if ppt_audience=="Internal Team" else "site actions, recovery plans, this week priorities"}
- Highlight critical path impacts clearly
- No percentages in bullet text — use plain language
- Each slide must tell a clear story, not just list facts"""

            raw = call_groq(prompt, max_tokens=2000)
            try:
                clean = raw.strip()
                if "```" in clean:
                    clean = clean.split("```")[1]
                    if clean.startswith("json"):
                        clean = clean[4:]
                parsed = json.loads(clean.strip())
                slides_data = parsed.get("slides", [])
            except Exception:
                slides_data = []
                st.error("Could not parse slide content. Try generating again.")

        # ── Preview slides ────────────────────────────────────
        if slides_data:
            st.divider()
            st.markdown("#### Preview — Review and Edit Before Downloading")

            edited_slides = []
            for i, slide in enumerate(slides_data):
                with st.expander(
                    f"Slide {i+2}: {slide.get('title','')}", expanded=True
                ):
                    col_preview, col_edit = st.columns([1, 1])
                    with col_preview:
                        st.markdown(f"**{slide.get('title','')}**")
                        for b in slide.get("bullets",[]):
                            st.markdown(f"▸ {b}")
                        if slide.get("highlight"):
                            st.info(f"📌 {slide['highlight']}")
                    with col_edit:
                        st.caption("Edit content:")
                        new_title = st.text_input(
                            "Title", value=slide.get("title",""),
                            key=f"edit_title_{i}"
                        )
                        new_bullets_raw = st.text_area(
                            "Bullets (one per line)",
                            value="\n".join(slide.get("bullets",[])),
                            height=120, key=f"edit_bullets_{i}"
                        )
                        new_highlight = st.text_input(
                            "Key stat", value=slide.get("highlight",""),
                            key=f"edit_highlight_{i}"
                        )
                        edited_slides.append({
                            "section":   slide.get("section",""),
                            "title":     new_title,
                            "bullets":   [b.strip() for b in new_bullets_raw.split("\n") if b.strip()],
                            "highlight": new_highlight,
                        })

            # ── Build graphs for embedding ────────────────────
            graph_images = {}
            if graph_selected and not history_df.empty:
                with st.spinner("Generating graphs for embedding..."):
                    if ppt_g_scurve:
                        scurve = history_df.groupby("week_number").agg(
                            actual=("cum_actual_pct","mean"),
                            planned=("cum_planned_pct","mean")
                        ).reset_index()
                        scurve[["actual","planned"]] *= 100
                        fig, ax = plt.subplots(figsize=(8,4))
                        ax.plot(scurve["week_number"], scurve["planned"],
                                color=NAVY, linewidth=2, label="Planned %",
                                linestyle="--", marker="o", markersize=3)
                        ax.plot(scurve["week_number"], scurve["actual"],
                                color=AMBER, linewidth=2, label="Actual %",
                                marker="o", markersize=3)
                        ax.fill_between(scurve["week_number"],
                                        scurve["planned"], scurve["actual"],
                                        where=scurve["actual"]<scurve["planned"],
                                        alpha=0.12, color=RED)
                        ax.set_title("S-Curve: Planned vs Actual", fontweight="bold")
                        ax.legend(fontsize=8); ax.grid(alpha=0.3)
                        ax.set_facecolor("#FAFAFA"); fig.patch.set_facecolor("white")
                        plt.tight_layout()
                        buf = io.BytesIO()
                        fig.savefig(buf, format="png", dpi=120, bbox_inches="tight")
                        buf.seek(0); graph_images["scurve"] = buf.getvalue()
                        plt.close(fig)

                    if ppt_g_variance:
                        trend = (history_df.groupby("week_number")["variance_pct"]
                                 .mean().reset_index())
                        trend["variance_pct"] *= 100
                        fig, ax = plt.subplots(figsize=(8,3.5))
                        colors = [RED if v < 0 else GREEN for v in trend["variance_pct"]]
                        ax.bar(trend["week_number"], trend["variance_pct"],
                               color=colors, width=0.6)
                        ax.axhline(0, color=GRAY, linewidth=0.8, linestyle="--")
                        ax.set_title("Average Variance by Week", fontweight="bold")
                        ax.grid(axis="y", alpha=0.3)
                        ax.set_facecolor("#FAFAFA"); fig.patch.set_facecolor("white")
                        plt.tight_layout()
                        buf = io.BytesIO()
                        fig.savefig(buf, format="png", dpi=120, bbox_inches="tight")
                        buf.seek(0); graph_images["variance"] = buf.getvalue()
                        plt.close(fig)

                    if ppt_g_phase and not activities.empty:
                        pd_data = activities.groupby("phase").agg(
                            actual=("cum_actual_pct","mean"),
                            planned=("cum_planned_pct","mean")
                        ).reset_index()
                        pd_data[["actual","planned"]] *= 100
                        fig, ax = plt.subplots(figsize=(8,4))
                        x = range(len(pd_data)); w = 0.35
                        ax.bar([i-w/2 for i in x], pd_data["planned"],
                               w, label="Planned %", color=NAVY, alpha=0.85)
                        ax.bar([i+w/2 for i in x], pd_data["actual"],
                               w, label="Actual %", color=AMBER, alpha=0.85)
                        ax.set_xticks(list(x))
                        ax.set_xticklabels(pd_data["phase"])
                        ax.set_title("Phase-wise Completion", fontweight="bold")
                        ax.legend(); ax.grid(axis="y", alpha=0.3)
                        ax.set_facecolor("#FAFAFA"); fig.patch.set_facecolor("white")
                        plt.tight_layout()
                        buf = io.BytesIO()
                        fig.savefig(buf, format="png", dpi=120, bbox_inches="tight")
                        buf.seek(0); graph_images["phase"] = buf.getvalue()
                        plt.close(fig)

                    if ppt_g_slip and not activities.empty:
                        slip_data = activities[activities["weeks_slip"]>0].sort_values(
                            "weeks_slip", ascending=True)
                        if not slip_data.empty:
                            fig, ax = plt.subplots(
                                figsize=(8, max(3, len(slip_data)*0.4)))
                            labels = [f"{r.get('phase','')} · {r.get('activity','')}"
                                      for _, r in slip_data.iterrows()]
                            colors = [RED if r.get("is_critical_path") else AMBER
                                      for _, r in slip_data.iterrows()]
                            ax.barh(labels, slip_data["weeks_slip"],
                                    color=colors, height=0.5)
                            ax.set_title("Weeks Slip by Activity", fontweight="bold")
                            ax.grid(axis="x", alpha=0.3)
                            ax.set_facecolor("#FAFAFA"); fig.patch.set_facecolor("white")
                            plt.tight_layout()
                            buf = io.BytesIO()
                            fig.savefig(buf, format="png", dpi=120, bbox_inches="tight")
                            buf.seek(0); graph_images["slip"] = buf.getvalue()
                            plt.close(fig)

            # ── Build PPTX ────────────────────────────────────
            try:
                from pptx import Presentation as PPTXPres
                from pptx.util import Inches, Pt, Emu
                from pptx.dml.color import RGBColor as RGB

                prs = PPTXPres()
                prs.slide_width  = Inches(13.33)
                prs.slide_height = Inches(7.5)
                blank = prs.slide_layouts[6]

                C_NAVY  = RGB(0x1E, 0x27, 0x61)
                C_WHITE = RGB(0xFF, 0xFF, 0xFF)
                C_AMBER = RGB(0xE8, 0x77, 0x22)
                C_GRAY  = RGB(0x44, 0x44, 0x44)
                C_RED   = RGB(0xD3, 0x2F, 0x2F)

                def pptx_slide(title_text, bullets, highlight=None,
                               is_cover=False, note_text=""):
                    sl = prs.slides.add_slide(blank)
                    # Background
                    bg = sl.shapes.add_shape(
                        1, 0, 0, prs.slide_width, prs.slide_height)
                    bg.fill.solid()
                    bg.fill.fore_color.rgb = C_NAVY if is_cover else C_WHITE
                    bg.line.fill.background()
                    # Top amber bar
                    bar = sl.shapes.add_shape(
                        1, 0, 0, prs.slide_width, Inches(0.08))
                    bar.fill.solid()
                    bar.fill.fore_color.rgb = C_AMBER
                    bar.line.fill.background()
                    # Title
                    tb = sl.shapes.add_textbox(
                        Inches(0.5), Inches(0.18), Inches(11.5), Inches(1.1))
                    p = tb.text_frame.paragraphs[0]
                    p.text = title_text
                    p.font.size = Pt(30 if is_cover else 24)
                    p.font.bold = True
                    p.font.color.rgb = C_WHITE if is_cover else C_NAVY
                    if is_cover:
                        sub = sl.shapes.add_textbox(
                            Inches(0.5), Inches(1.6), Inches(12), Inches(0.8))
                        sp = sub.text_frame.paragraphs[0]
                        sp.text = f"Week {selected_week:02d}  ·  {project_name}  ·  {ppt_audience}"
                        sp.font.size = Pt(16)
                        sp.font.color.rgb = C_AMBER
                        details = sl.shapes.add_textbox(
                            Inches(0.5), Inches(2.8), Inches(12), Inches(2))
                        dtf = details.text_frame
                        dtf.word_wrap = True
                        for i, line in enumerate([
                            f"Baseline: {baseline_ver}",
                            f"Contractor: {contractor}",
                            f"Prepared by: {ppt_by}",
                        ]):
                            dp = dtf.paragraphs[0] if i==0 else dtf.add_paragraph()
                            dp.text = line
                            dp.font.size = Pt(13)
                            dp.font.color.rgb = C_WHITE
                    else:
                        # Highlight box
                        if highlight:
                            hb = sl.shapes.add_shape(
                                1, Inches(10.3), Inches(1.4),
                                Inches(2.6), Inches(1.5))
                            hb.fill.solid()
                            hb.fill.fore_color.rgb = C_AMBER
                            hb.line.fill.background()
                            ht = hb.text_frame.paragraphs[0]
                            ht.text = str(highlight)
                            ht.font.size = Pt(18)
                            ht.font.bold = True
                            ht.font.color.rgb = C_WHITE
                            ht.alignment = 2  # center
                        # Bullets
                        bb = sl.shapes.add_textbox(
                            Inches(0.5), Inches(1.4), Inches(9.6), Inches(5.6))
                        btf = bb.text_frame
                        btf.word_wrap = True
                        for i, bullet in enumerate(bullets):
                            bp = btf.paragraphs[0] if i==0 else btf.add_paragraph()
                            bp.text = f"▸  {bullet}"
                            bp.font.size = Pt(14)
                            bp.font.color.rgb = C_GRAY
                            bp.space_after = Pt(7)
                        # Custom note
                        if note_text:
                            nb = sl.shapes.add_textbox(
                                Inches(0.5), Inches(6.6), Inches(12), Inches(0.5))
                            np_ = nb.text_frame.paragraphs[0]
                            np_.text = f"📌 {note_text}"
                            np_.font.size = Pt(10)
                            np_.font.color.rgb = C_AMBER
                            np_.font.italic = True
                    # Footer
                    ft = sl.shapes.add_textbox(
                        Inches(0.5), Inches(7.15), Inches(12), Inches(0.28))
                    fp = ft.text_frame.paragraphs[0]
                    fp.text = (f"SiteInsight  ·  {project_name}  ·  "
                               f"Week {selected_week:02d}  ·  CONFIDENTIAL")
                    fp.font.size = Pt(7)
                    fp.font.color.rgb = RGB(0xAA, 0xAA, 0xAA)

                # Cover slide
                pptx_slide(
                    f"Weekly Progress Report — Week {selected_week:02d}",
                    [], is_cover=True
                )

                # Content slides in order
                for slide in edited_slides:
                    sec = slide.get("section","")
                    note = ""
                    if "Executive" in sec:     note = exec_note
                    elif "Action" in sec:      note = actions_note_ppt
                    elif "Next" in sec:        note = next_note_ppt
                    pptx_slide(
                        slide["title"],
                        slide["bullets"],
                        slide.get("highlight",""),
                        note_text=note
                    )

                # Graph slide(s)
                if graph_images:
                    imgs = list(graph_images.values())
                    # Up to 2 graphs per slide
                    for slide_start in range(0, len(imgs), 2):
                        sl = prs.slides.add_slide(blank)
                        bg = sl.shapes.add_shape(
                            1, 0, 0, prs.slide_width, prs.slide_height)
                        bg.fill.solid()
                        bg.fill.fore_color.rgb = C_WHITE
                        bg.line.fill.background()
                        bar = sl.shapes.add_shape(
                            1, 0, 0, prs.slide_width, Inches(0.08))
                        bar.fill.solid()
                        bar.fill.fore_color.rgb = C_AMBER
                        bar.line.fill.background()
                        tb = sl.shapes.add_textbox(
                            Inches(0.5), Inches(0.18), Inches(12), Inches(0.9))
                        p = tb.text_frame.paragraphs[0]
                        p.text = "Project Analytics — Charts"
                        p.font.size = Pt(22); p.font.bold = True
                        p.font.color.rgb = C_NAVY
                        chunk = imgs[slide_start:slide_start+2]
                        positions = (
                            [(Inches(0.4), Inches(1.1), Inches(12.5), Inches(6.0))]
                            if len(chunk)==1 else
                            [(Inches(0.3), Inches(1.1), Inches(6.3), Inches(6.0)),
                             (Inches(6.8), Inches(1.1), Inches(6.3), Inches(6.0))]
                        )
                        for img_bytes, (l, t, w, h) in zip(chunk, positions):
                            img_buf = io.BytesIO(img_bytes)
                            sl.shapes.add_picture(img_buf, l, t, w, h)
                        ft = sl.shapes.add_textbox(
                            Inches(0.5), Inches(7.15), Inches(12), Inches(0.28))
                        fp = ft.text_frame.paragraphs[0]
                        fp.text = (f"SiteInsight  ·  {project_name}  ·  "
                                   f"Week {selected_week:02d}  ·  CONFIDENTIAL")
                        fp.font.size = Pt(7)
                        fp.font.color.rgb = RGB(0xAA, 0xAA, 0xAA)

                buf = io.BytesIO()
                prs.save(buf)
                buf.seek(0)
                st.divider()
                st.download_button(
                    "⬇️ Download PowerPoint (.pptx)", buf,
                    f"SiteInsight_Wk{selected_week:02d}_{ppt_audience.replace(' ','_')}.pptx",
                    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                )
                st.success(
                    f"✅ Presentation ready — {len(edited_slides)+1} content slides "
                    f"+ {len(list(range(0,len(graph_images),2)))} graph slide(s). "
                    f"Open in PowerPoint to further customise."
                )

            except ImportError:
                st.warning("python-pptx not installed. Run: pip install python-pptx")


# ════════════════════════════════════════════════════════════════
# SECTION: GENERATE GRAPHS
# ════════════════════════════════════════════════════════════════
elif nav == "📈 Generate Graphs":
    st.markdown("### 📈 Graph Generator")

    with st.form("graph_form"):
        gc1, gc2 = st.columns(2)
        with gc1:
            g_scurve   = st.checkbox("S-Curve (planned vs actual)",    value=True)
            g_variance = st.checkbox("Variance trend — all weeks",      value=True)
            g_phase    = st.checkbox("Phase-wise completion",           value=True)
        with gc2:
            g_slip     = st.checkbox("Weeks slip by activity",          value=True)
            g_delay    = st.checkbox("Delay reasons breakdown",         value=True)
            g_critical = st.checkbox("Critical path: planned vs actual",value=True)
            g_all      = st.checkbox("Generate ALL graphs",             value=False)
        gen_graphs = st.form_submit_button("Generate Graphs", type="primary")

    if gen_graphs:
        if g_all:
            g_scurve = g_variance = g_phase = g_slip = g_delay = g_critical = True

        st.divider()

        # 0 · S-Curve
        if g_scurve and not history_df.empty:
            st.markdown("#### S-Curve — Planned vs Actual Progress")
            scurve = history_df.groupby("week_number").agg(
                actual=("cum_actual_pct","mean"),
                planned=("cum_planned_pct","mean")
            ).reset_index()
            scurve[["actual","planned"]] *= 100

            fig, ax = plt.subplots(figsize=(11,4.5))
            ax.plot(scurve["week_number"], scurve["planned"],
                    color=NAVY, linewidth=2.5, label="Planned %",
                    linestyle="--", marker="o", markersize=3)
            ax.plot(scurve["week_number"], scurve["actual"],
                    color=AMBER, linewidth=2.5, label="Actual %",
                    marker="o", markersize=3)
            ax.fill_between(scurve["week_number"],
                            scurve["planned"], scurve["actual"],
                            where=scurve["actual"] < scurve["planned"],
                            alpha=0.12, color=RED, label="Variance gap")
            ax.axvline(selected_week, color=AMBER, linewidth=1.5,
                       linestyle=":", alpha=0.8, label=f"Wk{selected_week}")
            ax.set_xlabel("Week Number", fontsize=10)
            ax.set_ylabel("Cumulative Progress (%)", fontsize=10)
            ax.set_title(
                f"{project_name} · S-Curve: Planned vs Actual Progress",
                fontweight="bold", fontsize=11)
            ax.legend(fontsize=9)
            ax.grid(alpha=0.3)
            ax.set_facecolor("#FAFAFA")
            fig.patch.set_facecolor("white")
            plt.tight_layout()
            buf = io.BytesIO()
            fig.savefig(buf, format="png", dpi=150, bbox_inches="tight")
            buf.seek(0)
            st.image(buf, use_column_width=True)
            st.download_button("⬇️ Download S-Curve", buf,
                               f"scurve_Wk{selected_week:02d}.png",
                               "image/png", key="dl_sc")
            plt.close(fig); st.divider()

        # 1 · Variance trend
        if g_variance and not history_df.empty:
            st.markdown("#### Variance trend — all weeks")
            trend = (history_df.groupby("week_number")["variance_pct"]
                     .mean().reset_index())
            trend["variance_pct"] *= 100
            fig, ax = plt.subplots(figsize=(10,3.5))
            colors = [RED if v < 0 else GREEN for v in trend["variance_pct"]]
            ax.bar(trend["week_number"], trend["variance_pct"],
                   color=colors, width=0.6, zorder=3)
            ax.axhline(0, color=GRAY, linewidth=0.8, linestyle="--")
            ax.axvline(selected_week-0.5, color=AMBER, linewidth=1.5,
                       linestyle="--", label=f"Wk{selected_week}")
            ax.set_xlabel("Week"); ax.set_ylabel("Avg Variance (%)")
            ax.set_title(f"{project_name} · Average Variance by Week",
                         fontweight="bold")
            ax.legend(); ax.grid(axis="y", alpha=0.3)
            ax.set_facecolor("#FAFAFA"); fig.patch.set_facecolor("white")
            plt.tight_layout()
            buf = io.BytesIO()
            fig.savefig(buf, format="png", dpi=150, bbox_inches="tight")
            buf.seek(0)
            st.image(buf, use_column_width=True)
            st.download_button("⬇️ Download", buf,
                               f"variance_Wk{selected_week:02d}.png",
                               "image/png", key="dl_v")
            plt.close(fig); st.divider()

        # 2 · Phase-wise completion
        if g_phase and not activities.empty:
            st.markdown("#### Phase-wise completion — this week")
            pd_data = activities.groupby("phase").agg(
                actual=("cum_actual_pct","mean"),
                planned=("cum_planned_pct","mean")
            ).reset_index()
            pd_data[["actual","planned"]] *= 100
            fig, ax = plt.subplots(figsize=(9,4))
            x = range(len(pd_data)); w = 0.35
            b1 = ax.bar([i-w/2 for i in x], pd_data["planned"],
                        w, label="Planned %", color=NAVY, alpha=0.85)
            b2 = ax.bar([i+w/2 for i in x], pd_data["actual"],
                        w, label="Actual %", color=AMBER, alpha=0.85)
            ax.set_xticks(list(x))
            ax.set_xticklabels(pd_data["phase"])
            ax.set_ylabel("Completion (%)")
            ax.set_title(f"Phase-wise Status · Week {selected_week:02d}",
                         fontweight="bold")
            ax.legend(); ax.grid(axis="y", alpha=0.3)
            ax.set_facecolor("#FAFAFA"); fig.patch.set_facecolor("white")
            for b in list(b1)+list(b2):
                ax.text(b.get_x()+b.get_width()/2, b.get_height()+0.5,
                        f"{b.get_height():.0f}%",
                        ha="center", va="bottom", fontsize=8)
            plt.tight_layout()
            buf = io.BytesIO()
            fig.savefig(buf, format="png", dpi=150, bbox_inches="tight")
            buf.seek(0)
            st.image(buf, use_column_width=True)
            st.download_button("⬇️ Download", buf,
                               f"phase_Wk{selected_week:02d}.png",
                               "image/png", key="dl_p")
            plt.close(fig); st.divider()

        # 3 · Slip tracker
        if g_slip and not activities.empty:
            st.markdown("#### Weeks slip by activity")
            slip_data = activities[
                activities["weeks_slip"] > 0
            ].sort_values("weeks_slip", ascending=True)
            if slip_data.empty:
                st.success("No activities with weeks slip this week.")
            else:
                fig, ax = plt.subplots(
                    figsize=(9, max(3.5, len(slip_data)*0.45)))
                colors = [RED if row.get("is_critical_path") else AMBER
                          for _, row in slip_data.iterrows()]
                # Use phase · activity as label
                labels = [
                    f"{r.get('phase','')} · {r.get('activity','')}"
                    for _, r in slip_data.iterrows()
                ]
                bars = ax.barh(labels, slip_data["weeks_slip"],
                               color=colors, height=0.55)
                ax.set_xlabel("Weeks Slip")
                ax.set_title(f"Weeks Slip · Week {selected_week:02d}",
                             fontweight="bold")
                ax.grid(axis="x", alpha=0.3)
                ax.set_facecolor("#FAFAFA"); fig.patch.set_facecolor("white")
                from matplotlib.patches import Patch
                ax.legend(handles=[
                    Patch(color=RED,   label="Critical"),
                    Patch(color=AMBER, label="Non-critical")
                ])
                for b in bars:
                    ax.text(b.get_width()+0.05,
                            b.get_y()+b.get_height()/2,
                            f"{int(b.get_width())}w",
                            va="center", fontsize=8)
                plt.tight_layout()
                buf = io.BytesIO()
                fig.savefig(buf, format="png", dpi=150, bbox_inches="tight")
                buf.seek(0)
                st.image(buf, use_column_width=True)
                st.download_button("⬇️ Download", buf,
                                   f"slip_Wk{selected_week:02d}.png",
                                   "image/png", key="dl_s")
                plt.close(fig); st.divider()

        # 4 · Delay reasons
        if g_delay and not activities.empty:
            st.markdown("#### Delay reasons breakdown")
            delayed = activities[
                (activities["variance_pct"] < 0) &
                (activities["delay_reason"].notna()) &
                (activities["delay_reason"] != "No Delay")
            ]
            if delayed.empty:
                st.success("No delay reasons to display.")
            else:
                counts  = delayed["delay_reason"].value_counts()
                fig, ax = plt.subplots(figsize=(7,4.5))
                palette = [NAVY, AMBER, RED, AMBER2, GREEN,
                           "#7B61FF","#FF6B6B","#4ECDC4"]
                wedges, texts, autotexts = ax.pie(
                    counts.values, labels=counts.index,
                    autopct="%1.0f%%",
                    colors=palette[:len(counts)],
                    startangle=140, pctdistance=0.82,
                )
                for t in texts:      t.set_fontsize(9)
                for at in autotexts:
                    at.set_fontsize(8); at.set_color("white")
                    at.set_fontweight("bold")
                ax.set_title(f"Delay Reasons · Week {selected_week:02d}",
                             fontweight="bold")
                fig.patch.set_facecolor("white"); plt.tight_layout()
                buf = io.BytesIO()
                fig.savefig(buf, format="png", dpi=150, bbox_inches="tight")
                buf.seek(0)
                st.image(buf, use_column_width=True)
                st.download_button("⬇️ Download", buf,
                                   f"delay_Wk{selected_week:02d}.png",
                                   "image/png", key="dl_d")
                plt.close(fig); st.divider()

        # 5 · Critical path
        if (g_critical and not activities.empty
                and "is_critical_path" in activities.columns):
            st.markdown("#### Critical path: planned vs actual")
            cp = activities[activities["is_critical_path"] == True].copy()
            if cp.empty:
                st.info("No critical path activities found.")
            else:
                cp = cp.sort_values("variance_pct")
                labels = [
                    f"{r.get('phase','')} · {r.get('activity','')}"
                    for _, r in cp.iterrows()
                ]
                fig, ax = plt.subplots(
                    figsize=(9, max(3.5, len(cp)*0.55)))
                x = range(len(cp))
                ax.barh([i+0.2 for i in x], cp["cum_planned_pct"]*100,
                        height=0.35, label="Planned %",
                        color=NAVY, alpha=0.8)
                ax.barh([i-0.2 for i in x], cp["cum_actual_pct"]*100,
                        height=0.35, label="Actual %",
                        color=AMBER, alpha=0.8)
                ax.set_yticks(list(x))
                ax.set_yticklabels(labels, fontsize=8)
                ax.set_xlabel("Completion (%)")
                ax.set_title(f"Critical Path · Week {selected_week:02d}",
                             fontweight="bold")
                ax.legend(); ax.grid(axis="x", alpha=0.3)
                ax.set_facecolor("#FAFAFA"); fig.patch.set_facecolor("white")
                plt.tight_layout()
                buf = io.BytesIO()
                fig.savefig(buf, format="png", dpi=150, bbox_inches="tight")
                buf.seek(0)
                st.image(buf, use_column_width=True)
                st.download_button("⬇️ Download", buf,
                                   f"critical_Wk{selected_week:02d}.png",
                                   "image/png", key="dl_c")
                plt.close(fig)


# ════════════════════════════════════════════════════════════════
# SECTION: RECOVERY SIMULATOR
# ════════════════════════════════════════════════════════════════
elif nav == "🔄 Recovery Simulator":
    st.markdown("### 🔄 Recovery Simulator")
    st.caption(
        "Select a delayed activity and simulate catch-up scenarios "
        "based on different acceleration targets."
    )

    if activities.empty or history_df.empty:
        st.info("No activity data available. Load history first.")
    else:
        behind_acts = activities[activities["variance_pct"] < -0.05].copy()
        if behind_acts.empty:
            st.success("No delayed activities this week.")
        else:
            act_options = [
                f"{row['phase']} · {row['activity']} "
                f"({row['variance_pct']*100:.1f}%)"
                for _, row in behind_acts.iterrows()
            ]
            selected_act_label = st.selectbox(
                "Select delayed activity to simulate", act_options)
            selected_act_id = behind_acts.iloc[
                act_options.index(selected_act_label)]["act_id"]
            act_row = behind_acts[
                behind_acts["act_id"] == selected_act_id].iloc[0]

            current_pct   = float(act_row["cum_actual_pct"] or 0) if pd.notna(act_row["cum_actual_pct"]) else 0.0
            planned_pct   = float(act_row["cum_planned_pct"] or 0) if pd.notna(act_row["cum_planned_pct"]) else 0.0
            variance      = float(act_row["variance_pct"] or 0) if pd.notna(act_row["variance_pct"]) else 0.0
            weeks_slip    = int(act_row["weeks_slip"]) if pd.notna(act_row["weeks_slip"]) else 0
            delay_reason  = act_row.get("delay_reason","—") or "—"
            resp_person   = act_row.get("responsible_person","—") or "—"
            hist_velocity = compute_weekly_velocity(history_df, selected_act_id)
            if pd.isna(hist_velocity) or hist_velocity is None:
                hist_velocity = 0.02

            st.divider()
            col_s1, col_s2, col_s3, col_s4, col_s5 = st.columns(5)
            col_s1.metric("Current actual", f"{current_pct*100:.1f}%")
            col_s2.metric("Planned target", f"{planned_pct*100:.1f}%")
            col_s3.metric("Variance",       f"{variance*100:.1f}%",
                          delta_color="inverse")
            col_s4.metric("Weeks slip",     f"{weeks_slip} wks",
                          delta_color="inverse")
            col_s5.metric("Responsible",    resp_person)
            st.caption(
                f"Delay reason: **{delay_reason}** · "
                f"Historical rate: **{hist_velocity*100:.2f}%/wk**"
            )

            st.divider()
            sc1, sc2 = st.columns(2)
            with sc1:
                acceleration = st.slider(
                    "Target weekly progress rate (%)",
                    min_value=1.0, max_value=25.0,
                    value=max(float(hist_velocity*100*1.5), 5.0),
                    step=0.5,
                )
            with sc2:
                total_project_weeks = st.number_input(
                    "Total project weeks",
                    min_value=20, max_value=200, value=80)

            accel_rate        = acceleration / 100
            realistic_rate    = hist_velocity if hist_velocity > 0 else 0.02
            optimistic_rate   = accel_rate
            conservative_rate = accel_rate * 0.6

            proj_realistic    = compute_projected_finish(
                current_pct, selected_week, realistic_rate, total_project_weeks)
            proj_optimistic   = compute_projected_finish(
                current_pct, selected_week, optimistic_rate, total_project_weeks)
            proj_conservative = compute_projected_finish(
                current_pct, selected_week, conservative_rate, total_project_weeks)

            st.divider()
            st.markdown("**Projected finish week**")
            r1, r2, r3 = st.columns(3)
            r1.metric("Conservative",
                      f"Week {min(proj_conservative, total_project_weeks)}",
                      delta_color="inverse")
            r2.metric("Target (your rate)",
                      f"Week {min(proj_optimistic, total_project_weeks)}",
                      delta_color="normal")
            r3.metric("At current rate",
                      f"Week {min(proj_realistic, total_project_weeks)}",
                      delta_color="inverse")

            st.divider()
            future_weeks = list(range(
                selected_week,
                min(selected_week+30, total_project_weeks+1)))

            def project_curve(start_pct, rate, weeks):
                pcts = []; p = start_pct
                for _ in weeks:
                    p = min(p+rate, 1.0); pcts.append(p*100)
                return pcts

            fig, ax = plt.subplots(figsize=(10,4))
            ax.plot(future_weeks,
                    project_curve(current_pct, optimistic_rate, future_weeks),
                    color=GREEN, linewidth=2,
                    label=f"Target ({acceleration:.1f}%/wk)")
            ax.plot(future_weeks,
                    project_curve(current_pct, conservative_rate, future_weeks),
                    color=AMBER2, linewidth=1.5, linestyle="--",
                    label=f"Conservative ({acceleration*0.6:.1f}%/wk)")
            ax.plot(future_weeks,
                    project_curve(current_pct, realistic_rate, future_weeks),
                    color=RED, linewidth=1.5, linestyle=":",
                    label=f"Current rate ({hist_velocity*100:.1f}%/wk)")
            ax.axhline(100, color=GRAY, linewidth=0.8,
                       linestyle="--", label="100% complete")
            ax.axhline(planned_pct*100, color=NAVY, linewidth=0.8,
                       linestyle="-.",
                       label=f"Planned ({planned_pct*100:.0f}%)")
            ax.set_xlabel("Week"); ax.set_ylabel("Cumulative Progress (%)")
            act_label = f"{act_row.get('phase','')} · {act_row['activity']}"
            ax.set_title(f"Recovery Simulation · {act_label}",
                         fontweight="bold")
            ax.legend(fontsize=9); ax.grid(alpha=0.3)
            ax.set_facecolor("#FAFAFA"); fig.patch.set_facecolor("white")
            plt.tight_layout()
            buf = io.BytesIO()
            fig.savefig(buf, format="png", dpi=150, bbox_inches="tight")
            buf.seek(0)
            st.image(buf, use_column_width=True)
            st.download_button("⬇️ Download simulation chart", buf,
                               f"recovery_{selected_act_id}_Wk{selected_week:02d}.png",
                               "image/png", key="dl_sim")
            plt.close(fig)

            st.divider()
            if st.button("Generate AI recovery recommendation", type="primary"):
                with st.spinner("Analysing via Groq..."):
                    prompt = f"""You are a senior construction planning consultant.
ACTIVITY: {act_row['activity']} ({selected_act_id})
PHASE: {act_row['phase']} | RESPONSIBLE: {resp_person}
CURRENT: {current_pct*100:.1f}% | PLANNED: {planned_pct*100:.1f}%
VARIANCE: {variance*100:.1f}% | WEEKS SLIP: {weeks_slip}
DELAY REASON: {delay_reason}
HISTORICAL RATE: {hist_velocity*100:.2f}%/wk | TARGET: {acceleration:.1f}%/wk
At current rate: Week {proj_realistic} | At target: Week {proj_optimistic}

Write a specific recovery recommendation (max 200 words):
1. Is {acceleration:.1f}%/week realistic given the delay reason?
2. What must {resp_person} do THIS WEEK specifically?
3. Consequence if not recovered by Week {selected_week+1}?
4. One direct instruction to the contractor.
Direct construction language only."""
                    rec = call_groq(prompt, max_tokens=400)
                st.markdown("**AI Recovery Recommendation**")
                st.warning(rec)


# ════════════════════════════════════════════════════════════════
# SECTION: EARLY WARNING
# ════════════════════════════════════════════════════════════════
elif nav == "⚠️ Early Warning":
    st.markdown("### ⚠️ Early Warning System")
    st.caption(
        "Identifies activities likely to slip NEXT week based on "
        "trend velocity — before it appears in the WPR."
    )

    LOOKBACK = 4

    if history_df.empty or activities.empty:
        st.info("Need at least 3 weeks of history for trend analysis.")
    else:
        warnings_list = []
        for _, act in activities.iterrows():
            aid      = act["act_id"]
            variance = float(act.get("variance_pct", 0) or 0)
            critical = bool(act.get("is_critical_path", False))
            slip     = int(act.get("weeks_slip", 0) or 0)
            resp     = act.get("responsible_person","—") or "—"
            slope    = compute_trend_slope(history_df, aid, LOOKBACK)
            risk_score = 0; risk_flags = []
            if slope < -0.015:
                risk_score += 3
                risk_flags.append(
                    f"Variance worsening {abs(slope)*100:.1f}%/wk over {LOOKBACK} wks")
            if variance < -0.05:
                risk_score += 2
                risk_flags.append(f"Already {abs(variance)*100:.1f}% behind plan")
            if slip >= 2:
                risk_score += 2
                risk_flags.append(f"{slip} weeks slip accumulated")
            if critical:
                risk_score += 2
                risk_flags.append("On critical path")
            if slope < 0 and variance < -0.03:
                risk_score += 1
                risk_flags.append("Negative trend + negative variance")
            if risk_score >= 3:
                risk_level = ("HIGH"   if risk_score >= 7 else
                              "MEDIUM" if risk_score >= 4 else "LOW")
                warnings_list.append({
                    "act_id":      aid,
                    "activity":    act["activity"],
                    "phase":       act["phase"],
                    "variance":    variance,
                    "slope":       slope,
                    "slip":        slip,
                    "critical":    critical,
                    "responsible": resp,
                    "risk_score":  risk_score,
                    "risk_level":  risk_level,
                    "flags":       risk_flags,
                })

        warnings_list.sort(key=lambda x: x["risk_score"], reverse=True)

        if not warnings_list:
            st.success("No early warning signals detected.")
        else:
            high   = [w for w in warnings_list if w["risk_level"]=="HIGH"]
            medium = [w for w in warnings_list if w["risk_level"]=="MEDIUM"]
            low    = [w for w in warnings_list if w["risk_level"]=="LOW"]

            w1, w2, w3 = st.columns(3)
            w1.metric("HIGH risk",   len(high),
                      delta="Immediate action" if high else None,
                      delta_color="inverse")
            w2.metric("MEDIUM risk", len(medium),
                      delta="Monitor closely" if medium else None,
                      delta_color="inverse")
            w3.metric("LOW risk",    len(low))
            st.divider()

            level_icons = {"HIGH":"🔴","MEDIUM":"🟡","LOW":"🟢"}
            for w in warnings_list:
                with st.expander(
                    f"{level_icons[w['risk_level']]} {w['risk_level']} · "
                    f"{w['phase']} · {w['activity']} "
                    f"(Score: {w['risk_score']})",
                    expanded=(w["risk_level"]=="HIGH"),
                ):
                    e1, e2, e3, e4, e5 = st.columns(5)
                    e1.metric("Variance",
                              f"{w['variance']*100:.1f}%",
                              delta_color="inverse")
                    e2.metric("Trend slope",
                              f"{w['slope']*100:.2f}%/wk",
                              delta_color=(
                                  "inverse" if w["slope"] < 0 else "normal"))
                    e3.metric("Weeks slip",
                              f"{w['slip']} wks",
                              delta_color=(
                                  "inverse" if w["slip"] > 0 else "normal"))
                    e4.metric("Critical",
                              "YES" if w["critical"] else "No")
                    e5.metric("Responsible", w["responsible"])
                    st.markdown("**Why flagged:**")
                    for flag in w["flags"]:
                        st.markdown(f"- {flag}")

            # Trend chart
            st.divider()
            st.markdown("#### Variance trend — at-risk activities")
            at_risk_ids = [w["act_id"] for w in warnings_list[:6]]
            hist_subset = history_df[history_df["act_id"].isin(at_risk_ids)]
            if not hist_subset.empty:
                fig, ax = plt.subplots(figsize=(10,4))
                for aid in at_risk_ids:
                    act_data = hist_subset[
                        hist_subset["act_id"]==aid
                    ].sort_values("week_number")
                    if len(act_data) < 2:
                        continue
                    risk_level = next(
                        (w["risk_level"] for w in warnings_list
                         if w["act_id"]==aid), "LOW")
                    color = (RED   if risk_level=="HIGH" else
                             AMBER if risk_level=="MEDIUM" else GRAY)
                    # Use phase · activity as legend label
                    w_item = next(
                        (w for w in warnings_list if w["act_id"]==aid), {})
                    leg_label = (
                        f"{w_item.get('phase','')} · "
                        f"{w_item.get('activity','')}"
                    )
                    ax.plot(act_data["week_number"],
                            act_data["variance_pct"]*100,
                            label=leg_label,
                            linewidth=2.5 if risk_level=="HIGH" else 1.5,
                            color=color, marker="o", markersize=3)
                ax.axhline(0,   color=GRAY, linewidth=0.8, linestyle="--")
                ax.axhline(-5,  color=RED,  linewidth=0.5,
                           linestyle=":", alpha=0.5)
                ax.axhline(-10, color=RED,  linewidth=0.8,
                           linestyle=":", alpha=0.7)
                ax.axvline(selected_week, color=AMBER,
                           linewidth=1, linestyle="--", alpha=0.7)
                ax.set_xlabel("Week"); ax.set_ylabel("Variance (%)")
                ax.set_title("Variance Trajectory — At-risk Activities",
                             fontweight="bold")
                ax.legend(fontsize=7, ncol=2)
                ax.grid(alpha=0.3); ax.set_facecolor("#FAFAFA")
                fig.patch.set_facecolor("white"); plt.tight_layout()
                buf = io.BytesIO()
                fig.savefig(buf, format="png", dpi=150, bbox_inches="tight")
                buf.seek(0)
                st.image(buf, use_column_width=True)
                st.download_button("⬇️ Download early warning chart", buf,
                                   f"early_warning_Wk{selected_week:02d}.png",
                                   "image/png", key="dl_ew")
                plt.close(fig)

            st.divider()
            if st.button("Generate proactive weekly brief", type="primary"):
                with st.spinner("Generating via Groq..."):
                    high_names = [
                        f"{w['phase']} · {w['activity']} "
                        f"(responsible: {w['responsible']}, "
                        f"slope {w['slope']*100:.2f}%/wk)"
                        for w in high
                    ]
                    med_names = [
                        f"{w['phase']} · {w['activity']} "
                        f"(responsible: {w['responsible']})"
                        for w in medium
                    ]
                    prompt = f"""You are a proactive construction project intelligence system.
Week {selected_week} of {project_name}.

EARLY WARNING ({LOOKBACK}-week trend analysis):
HIGH RISK — will worsen next week: {high_names if high_names else 'None'}
MEDIUM RISK — monitor: {med_names if med_names else 'None'}

Write a PROACTIVE brief (max 250 words):
1. Predict what happens in Week {selected_week+1} if no action taken
2. Name each responsible person and their specific required action THIS WEEK
3. Identify milestone / critical path impact
4. End with ONE thing the PM must do before next site meeting

This is intelligence — tell the team what WILL happen, not what HAS happened.
Direct, action-oriented construction language."""
                    brief = call_groq(prompt, max_tokens=500)
                st.markdown("**Proactive Weekly Brief — Week ahead**")
                st.warning(brief)
