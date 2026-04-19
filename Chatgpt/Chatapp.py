"""
app.py
──────
SiteInsight · Streamlit Dashboard — Final Version
Tabs:
  1. Intelligence Report  (+ responsible person, flag status, site photos)
  2. Generate MOM         (+ responsible person, email notifications)
  3. Generate PPT
  4. Generate Graphs
  5. Recovery Simulator
  6. Early Warning

Run:
    python -m streamlit run app.py
"""

import os
import io
import json
import smtplib
import pandas as pd
import streamlit as st
import matplotlib.pyplot as plt
import matplotlib
matplotlib.use("Agg")

from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
from supabase import create_client, Client
from groq import Groq
from dotenv import load_dotenv
from datetime import datetime, date

load_dotenv()

SUPABASE_URL   = os.environ["SUPABASE_URL"]
SUPABASE_KEY   = os.environ["SUPABASE_KEY"]
GROQ_API_KEY   = os.environ["GROQ_API_KEY"]
EMAIL_SENDER   = os.environ.get("EMAIL_SENDER", "")
EMAIL_PASSWORD = os.environ.get("EMAIL_PASSWORD", "")
GROQ_MODEL     = "llama-3.3-70b-versatile"

st.set_page_config(
    page_title="SiteInsight · WPR Intelligence",
    page_icon="🏗",
    layout="wide",
)

# ── Clients ───────────────────────────────────────────────────
@st.cache_resource
def get_supabase() -> Client:
    return create_client(SUPABASE_URL, SUPABASE_KEY)

@st.cache_resource
def get_groq():
    return Groq(api_key=GROQ_API_KEY)

supabase    = get_supabase()
groq_client = get_groq()

# ── Placeholder images per activity type ──────────────────────
PHOTO_PLACEHOLDERS = {
    "excavation":    ("https://images.unsplash.com/photo-1581094794329-c8112a89af12?w=400",
                      "Excavation works in progress"),
    "shuttering":    ("https://images.unsplash.com/photo-1504307651254-35680f356dfd?w=400",
                      "Formwork and shuttering"),
    "reinforcement": ("https://images.unsplash.com/photo-1590579491624-f98f36d4c763?w=400",
                      "Reinforcement steel works"),
    "concrete":      ("https://images.unsplash.com/photo-1503387762-592deb58ef4e?w=400",
                      "Concrete pouring and casting"),
    "waterproofing": ("https://images.unsplash.com/photo-1558618666-fcd25c85cd64?w=400",
                      "Waterproofing membrane application"),
    "default":       ("https://images.unsplash.com/photo-1486325212027-8081e485255e?w=400",
                      "Construction site progress"),
}

def get_photo_for_activity(activity_name: str) -> tuple:
    name = activity_name.lower()
    for key in PHOTO_PLACEHOLDERS:
        if key in name:
            return PHOTO_PLACEHOLDERS[key]
    return PHOTO_PLACEHOLDERS["default"]


# ── Data loaders ──────────────────────────────────────────────
@st.cache_data(ttl=60)
def load_weeks():
    r = supabase.table("wpr_headers").select(
        "week_number,load_type").order("week_number").execute()
    return r.data or []

@st.cache_data(ttl=60)
def load_header(week):
    r = supabase.table("wpr_headers").select("*").eq("week_number", week).execute()
    return r.data[0] if r.data else {}

@st.cache_data(ttl=60)
def load_activities(week):
    r = supabase.table("wpr_activities").select("*").eq(
        "week_number", week).order("act_id").execute()
    return pd.DataFrame(r.data) if r.data else pd.DataFrame()

@st.cache_data(ttl=30)
def load_detections(week):
    r = supabase.table("detection_results").select("*").eq(
        "week_number", week).order("rule_id").execute()
    return pd.DataFrame(r.data) if r.data else pd.DataFrame()

@st.cache_data(ttl=60)
def load_narrative(week):
    r = supabase.table("ai_narratives").select("*").eq("week_number", week).execute()
    return r.data[0] if r.data else {}

@st.cache_data(ttl=60)
def load_dq_flags(week):
    r = supabase.table("dq_flags").select("*").eq("week_number", week).execute()
    return pd.DataFrame(r.data) if r.data else pd.DataFrame()

@st.cache_data(ttl=300)
def load_all_history():
    r = supabase.table("wpr_activities").select(
        "week_number,act_id,phase,activity,variance_pct,weeks_slip,"
        "cum_actual_pct,cum_planned_pct,delay_reason,is_critical_path,"
        "total_scope,responsible_person"
    ).order("week_number").execute()
    return pd.DataFrame(r.data) if r.data else pd.DataFrame()

@st.cache_data(ttl=300)
def load_baselines():
    r = supabase.table("baselines").select("*").order("id").execute()
    return pd.DataFrame(r.data) if r.data else pd.DataFrame()


# ── Groq helper ───────────────────────────────────────────────
def call_groq(prompt: str, max_tokens: int = 1500) -> str:
    response = groq_client.chat.completions.create(
        model=GROQ_MODEL,
        messages=[{"role": "user", "content": prompt}],
        temperature=0.3,
        max_tokens=max_tokens,
    )
    return response.choices[0].message.content.strip()


# ── Email helper ──────────────────────────────────────────────
def send_email(to_email: str, subject: str, body: str) -> bool:
    try:
        msg = MIMEMultipart("alternative")
        msg["Subject"] = subject
        msg["From"]    = EMAIL_SENDER
        msg["To"]      = to_email
        msg.attach(MIMEText(body, "plain"))
        with smtplib.SMTP_SSL("smtp.gmail.com", 465) as server:
            server.login(EMAIL_SENDER, EMAIL_PASSWORD)
            server.sendmail(EMAIL_SENDER, to_email, msg.as_string())
        return True
    except Exception as e:
        st.error(f"Email failed to {to_email}: {e}")
        return False


# ── Helper functions ──────────────────────────────────────────
def compute_weekly_velocity(hdf: pd.DataFrame, aid: str) -> float:
    act_hist = hdf[hdf["act_id"] == aid].sort_values("week_number")
    if len(act_hist) < 2:
        return 0.0
    return float(act_hist["cum_actual_pct"].diff().dropna().mean())

def compute_projected_finish(current_pct, current_week, weekly_rate, total_weeks=80):
    if weekly_rate <= 0:
        return total_weeks
    remaining = 1.0 - current_pct
    return int(current_week + remaining / weekly_rate)

def compute_trend_slope(hdf: pd.DataFrame, aid: str, lookback: int = 4) -> float:
    act_hist = hdf[hdf["act_id"] == aid].sort_values("week_number").tail(lookback)
    if len(act_hist) < 2:
        return 0.0
    variances = act_hist["variance_pct"].values
    return float(variances[-1] - variances[0]) / len(variances)

def compute_risk(act, history_df, lookback=4):
    aid = act["act_id"]
    slope = compute_trend_slope(history_df, aid, lookback)

    variance = float(act.get("variance_pct", 0) or 0)
    slip = int(act.get("weeks_slip", 0) or 0)
    critical = bool(act.get("is_critical_path", False))

    score = 0
    flags = []

    if slope < -0.015:
        score += 3
        flags.append(f"Worsening trend ({abs(slope)*100:.1f}%/wk)")

    if variance < -0.05:
        score += 2
        flags.append(f"{abs(variance)*100:.1f}% behind plan")

    if slip >= 2:
        score += 2
        flags.append(f"{slip} weeks slip")

    if critical:
        score += 2
        flags.append("Critical path")

    if slope < 0 and variance < -0.03:
        score += 1
        flags.append("Negative trend + delay")

    return score, flags, slope

def update_flag_status(flag_id: int, new_status: str):
    supabase.table("detection_results").update(
        {"status": new_status}
    ).eq("id", flag_id).execute()


# ── Chart colours ─────────────────────────────────────────────
NAVY   = "#1E2761"
ACCENT = "#00A8E8"
RED    = "#E24B4A"
AMBER  = "#F9A825"
GREEN  = "#2E7D32"
GRAY   = "#888780"


# ════════════════════════════════════════════════════════════════
# SIDEBAR
# ════════════════════════════════════════════════════════════════
with st.sidebar:
    st.title("🏗 SiteInsight")
    st.caption("WPR Intelligence · Horizon Commercial Tower")
    st.divider()

    weeks_data = load_weeks()
    if not weeks_data:
        st.error("No data loaded yet. Run load_history.py first.")
        st.stop()

    week_options = [w["week_number"] for w in weeks_data]
    live_weeks   = [w["week_number"] for w in weeks_data
                    if w.get("load_type") == "live"]

    selected_week = st.selectbox(
        "Select week",
        options=week_options,
        index=len(week_options) - 1,
        format_func=lambda w:
            f"Week {w:02d} {'🟢 Live' if w in live_weeks else '📁 History'}"
    )

    st.divider()
    baselines_df = load_baselines()
    if not baselines_df.empty:
        st.markdown("**Baselines**")
        for _, row in baselines_df.iterrows():
            active = "✅" if row["is_active"] else ""
            st.caption(
                f"{active} **{row['baseline_version']}** · "
                f"ends {row['planned_end_date']}"
            )

    st.divider()
    if st.button("🔄 Refresh data"):
        st.cache_data.clear()
        st.rerun()


# ── Load data ─────────────────────────────────────────────────

header     = load_header(selected_week)
activities = load_activities(selected_week)
detections = load_detections(selected_week)
narrative  = load_narrative(selected_week)
dq_flags   = load_dq_flags(selected_week)
history_df = load_all_history()
# ── Risk Engine (NEW) ─────────────────────────────────────────

risk_list = []

if not activities.empty and not history_df.empty:
    for _, act in activities.iterrows():
        score, flags, slope = compute_risk(act, history_df)

        if score > 0:
            risk_list.append({
                "act_id": act["act_id"],
                "activity": act["activity"],
                "score": score,
                "flags": flags,
                "slope": slope,
                "variance": act["variance_pct"],
                "critical": act.get("is_critical_path", False),
                "responsible": act.get("responsible_person", "—")
            })

risk_df = pd.DataFrame(risk_list)

if not risk_df.empty:
    risk_df = risk_df.sort_values("score", ascending=False)

project_name = header.get("project_name", "Horizon Commercial Tower")
baseline_ver = header.get("baseline_version", "B1")
contractor   = header.get("contractor", "—")
report_date  = header.get("report_date", "—")

# ── Page header ───────────────────────────────────────────────
st.title(f"Week {selected_week:02d} · WPR Intelligence Report")
c1, c2, c3, c4 = st.columns(4)
c1.caption(f"**Project:** {project_name}")
c2.caption(f"**Baseline:** {baseline_ver}")
c3.caption(f"**Contractor:** {contractor}")
c4.caption(f"**Report date:** {report_date}")

if header.get("baseline_revised_this_week"):
    st.warning(
        f"⚠️ Baseline revised this week · "
        f"{header.get('baseline_revision_note','')}"
    )

st.divider()


# ════════════════════════════════════════════════════════════════
# TABS
# ════════════════════════════════════════════════════════════════
tab1, tab2, tab3, tab4, tab5, tab6 = st.tabs([
    "📋 Intelligence Report",
    "📝 Generate MOM",
    "📊 Generate PPT",
    "📈 Generate Graphs",
    "🔄 Recovery Simulator",
    "⚠️ Early Warning",
])


# ════════════════════════════════════════════════════════════════
# TAB 1 · INTELLIGENCE REPORT
# ════════════════════════════════════════════════════════════════
with tab1:
    # ── Top Risks (NEW) ────────────────────────────────────────

 if not risk_df.empty:
    st.subheader("🚨 Top Risks This Week")

    top_risks = risk_df.sort_values("score", ascending=False).head(3)

    for _, r in top_risks.iterrows():
        st.error(
            f"**{r['act_id']} — {r['activity']}**\n"
            f"→ {', '.join(r['flags'])}\n"
            f"👤 Responsible: {r['responsible']}"
        )

    st.divider()

    # ── KPIs ─────────────────────────────────────────────────
    if not activities.empty:
        total_acts = len(activities)
        behind     = int((activities["variance_pct"] < -0.05).sum())
        critical_behind = 0
        if "is_critical_path" in activities.columns:
            critical_behind = int((
                activities[activities["is_critical_path"] == True]
                ["variance_pct"] < -0.05
            ).sum())
        avg_variance = activities["variance_pct"].mean()
        max_slip     = int(activities["weeks_slip"].max()) \
                       if "weeks_slip" in activities.columns else 0

        k1, k2, k3, k4, k5 = st.columns(5)
        k1.metric("Total activities", total_acts)
        k2.metric("Behind plan", behind,
                  delta=f"{behind/total_acts*100:.0f}%",
                  delta_color="inverse")
        k3.metric("Critical path risk", critical_behind,
                  delta="HIGH" if critical_behind > 0 else "Clear",
                  delta_color="inverse" if critical_behind > 0 else "normal")
        k4.metric("Avg variance", f"{avg_variance*100:.1f}%",
                  delta_color="inverse")
        k5.metric("Max weeks slip", f"{max_slip} wks",
                  delta_color="inverse" if max_slip > 0 else "normal")

    st.divider()

    # ── AI Narrative ─────────────────────────────────────────
    
    st.subheader("📋 Executive Summary")
    if not activities.empty and not risk_df.empty:
           total_acts = len(activities)
           behind = int((activities["variance_pct"] < -0.05).sum())

           high_risk = risk_df[risk_df["score"] >= 7]

            # 👉 ADD THIS LINE
           top_issue = risk_df.sort_values("score", ascending=False).iloc[0]

           summary_text = f"""
         • {behind}/{total_acts} activities behind plan  
         • Top concern: {top_issue['act_id']} (Score: {top_issue['score']})  
         • {len(high_risk)} high-risk activities identified  
         """

           st.info(summary_text)
    else:
           st.info("No sufficient data for summary")

    col_r, col_rec = st.columns(2)
    with col_r:
            st.subheader("⚠️ Key Risks")
            st.markdown(narrative.get("key_risks", "—"))
    with col_rec:
            st.subheader("✅ Recommendations")
            st.markdown(narrative.get("recommendations", "—"))

    with st.expander("ℹ️ Model info"):
            st.caption(
                f"Model: {narrative.get('model_used','—')} · "
                f"Tokens: {narrative.get('prompt_tokens',0)} prompt + "
                f"{narrative.get('completion_tokens',0)} completion"
            )

    # Trend commentary
    if not history_df.empty:
            st.subheader("📉 Trend Analysis")
            wk_trend = (history_df.groupby("week_number")["variance_pct"]
                        .mean().reset_index())
            last3 = wk_trend.tail(3)["variance_pct"].values
            if len(last3) >= 2:
                direction = last3[-1] - last3[0]
                if direction < -0.01:
                    st.error(
                        f"Project variance is **worsening** — dropped "
                        f"{abs(direction)*100:.1f}% over last 3 weeks. "
                        f"Intervention required."
                    )
                elif direction > 0.01:
                    st.success(
                        f"Project variance is **recovering** — improved "
                        f"{direction*100:.1f}% over last 3 weeks. "
                        f"Maintain recovery momentum."
                    )
                else:
                    st.warning(
                        "Project variance is **stable but negative** — "
                        "no recovery visible yet."
                    )
    else:
        st.info("No AI narrative for this week. Run pipeline.py to generate.")

    st.divider()

    # ── Detection Flags with Status + Responsible Person ─────
    st.subheader(f"🔍 Detection Flags ({len(detections)} this week)")
    rule_names    = {
        "R1": "Schedule Slip", "R2": "Scope Creep",
        "R3": "Stale Issue",   "R4": "Critical Path at Risk"
    }
    severity_icon = {"high": "🔴", "medium": "🟡", "low": "🟢"}
    status_options = ["Open", "In Progress", "Resolved"]

    if detections.empty:
        st.success("No flags raised this week.")
    else:
        for rule_id, rname in rule_names.items():
            hits = detections[detections["rule_id"] == rule_id]
            if hits.empty:
                continue
            with st.expander(
                f"{rule_id} · {rname} "
                f"({len(hits)} hit{'s' if len(hits)>1 else ''})",
                expanded=(rule_id in ["R1","R4"])
            ):
                for _, h in hits.iterrows():
                    icon    = severity_icon.get(h.get("severity","low"), "⚪")
                    act_id  = h["act_id"]
                    flag_id = h.get("id")
                    current_status = h.get("status", "Open") or "Open"

                    # Responsible person from activities
                    resp_person = "—"
                    if not activities.empty and "responsible_person" in activities.columns:
                        match = activities[activities["act_id"] == act_id]
                        if not match.empty:
                            rp = match.iloc[0].get("responsible_person")
                            if rp and str(rp) not in ["nan","None",""]:
                                resp_person = str(rp)

                    fc1, fc2, fc3 = st.columns([3, 1, 1])
                    with fc1:
                        st.markdown(
                            f"{icon} **{act_id}** — {h['detail']}\n\n"
                            f"👤 Responsible: **{resp_person}**"
                        )
                    with fc2:
                        st.caption("Status")
                        new_status = st.selectbox(
                            "status",
                            status_options,
                            index=status_options.index(current_status),
                            key=f"status_{flag_id}_{rule_id}_{act_id}",
                            label_visibility="collapsed",
                        )
                        if new_status != current_status and flag_id:
                            update_flag_status(flag_id, new_status)
                            st.cache_data.clear()
                    with fc3:
                        if current_status == "Resolved":
                            st.success("✓ Resolved")
                        elif current_status == "In Progress":
                            st.warning("⏳ In Progress")
                        else:
                            st.error("🔴 Open")

                    st.divider()

    # ── Flag summary from last week ───────────────────────────
    if selected_week > 1 and not detections.empty:
        prev_week = selected_week - 1
        prev_det  = supabase.table("detection_results").select(
            "status"
        ).eq("week_number", prev_week).execute()
        if prev_det.data:
            prev_df     = pd.DataFrame(prev_det.data)
            total_prev  = len(prev_df)
            resolved    = (prev_df["status"] == "Resolved").sum()
            in_progress = (prev_df["status"] == "In Progress").sum()
            still_open  = (prev_df["status"] == "Open").sum()
            st.info(
                f"**From Week {prev_week:02d}:** {total_prev} flags raised · "
                f"✓ {resolved} resolved · "
                f"⏳ {in_progress} in progress · "
                f"🔴 {still_open} still open"
            )

    st.divider()

    # ── Activity Table ────────────────────────────────────────
    st.subheader("📊 Activity Progress")
    if not activities.empty:
        phases    = ["All"] + sorted(
            activities["phase"].dropna().unique().tolist())
        sel_phase = st.selectbox(
            "Filter by phase", phases, key="phase_tab1")
        disp      = activities.copy()
        if sel_phase != "All":
            disp = disp[disp["phase"] == sel_phase]

        show_cols = [c for c in [
            "act_id","phase","activity","unit",
            "cum_actual_pct","cum_planned_pct","variance_pct",
            "weeks_slip","delay_reason","responsible_person",
            "is_critical_path","remarks"
        ] if c in disp.columns]
        disp = disp[show_cols].copy()

        for pct_col in ["cum_actual_pct","cum_planned_pct","variance_pct"]:
            if pct_col in disp.columns:
                disp[pct_col] = (disp[pct_col]*100).round(1).astype(str)+"%"
        if "is_critical_path" in disp.columns:
            disp["is_critical_path"] = disp["is_critical_path"].map(
                {True:"✅", False:"—"})

        disp.columns = [c.replace("_"," ").title() for c in disp.columns]
        st.dataframe(disp, use_container_width=True, hide_index=True)

    st.divider()

    # ── Site Photos Gallery ───────────────────────────────────
    st.subheader("📸 Site Photos")
    st.caption(
        "Illustrative photos matched to activity type. "
        "In production, site engineer uploads actual photos."
    )

    if not activities.empty:
        phases_list = activities["phase"].dropna().unique().tolist()
        for phase in phases_list:
            phase_acts = activities[activities["phase"] == phase]
            st.markdown(f"**{phase}**")
            photo_cols = st.columns(min(len(phase_acts), 3))
            for idx, (_, act_row) in enumerate(
                phase_acts.head(3).iterrows()
            ):
                img_url, caption = get_photo_for_activity(
                    act_row.get("activity",""))
                with photo_cols[idx % 3]:
                    try:
                        st.image(
                            img_url,
                            caption=(
                                f"{act_row.get('act_id','')} · "
                                f"{caption} · "
                                f"Actual: "
                                f"{act_row.get('cum_actual_pct',0)*100:.0f}%"
                            ),
                            use_column_width=True,
                        )
                    except Exception:
                        st.info(f"📷 {act_row.get('activity','')} — photo unavailable")
            st.divider()


# ════════════════════════════════════════════════════════════════
# TAB 2 · GENERATE MOM
# ════════════════════════════════════════════════════════════════
with tab2:
    st.subheader("📝 Minutes of Meeting Generator")
    st.caption(
        "Auto-generates a structured MOM with responsible persons. "
        "Send personalised email to each team member."
    )

    with st.form("mom_form"):
        st.markdown("**Meeting details**")
        mc1, mc2 = st.columns(2)
        meeting_date        = mc1.date_input("Meeting date", value=date.today())
        meeting_chairperson = mc2.text_input(
            "Chaired by", value="Project Manager")

        attendees = st.text_area(
            "Attendees (one per line)",
            value=(
                "Project Manager\nPlanning Engineer\n"
                "Site Engineer\nContractor Representative"
            ),
            height=100,
        )

        st.markdown("**Sections to include**")
        sc1, sc2, sc3 = st.columns(3)
        inc_delays   = sc1.checkbox(
            "Delayed activities & accountabilities", value=True)
        inc_critical = sc2.checkbox("Critical path risks", value=True)
        inc_stale    = sc3.checkbox("Stale / recurring issues", value=True)
        sc4, sc5, sc6 = st.columns(3)
        inc_actions  = sc4.checkbox("Action items with owners", value=True)
        inc_baseline = sc5.checkbox("Baseline status", value=True)
        inc_next     = sc6.checkbox("Next week targets", value=True)

        gen_mom = st.form_submit_button("Generate MOM", type="primary")

    if gen_mom:
        if activities.empty:
            st.warning("No activity data for this week.")
        else:
            with st.spinner("Generating MOM via Groq..."):
                behind_acts   = activities[
                    activities["variance_pct"] < -0.05]
                critical_acts = activities[
                    (activities.get("is_critical_path","") == True) &
                    (activities["variance_pct"] < -0.05)
                ] if "is_critical_path" in activities.columns \
                    else pd.DataFrame()
                det_r3 = detections[
                    detections["rule_id"] == "R3"
                ] if not detections.empty else pd.DataFrame()

                sections = (
                    (["delayed_activities"] if inc_delays   else []) +
                    (["critical_path"]      if inc_critical else []) +
                    (["stale_issues"]       if inc_stale    else []) +
                    (["action_items"]       if inc_actions  else []) +
                    (["baseline_status"]    if inc_baseline else []) +
                    (["next_week_targets"]  if inc_next     else [])
                )

                # Build behind activities string with responsible person
                behind_str = "None"
                if not behind_acts.empty:
                    rows = []
                    for _, r in behind_acts.iterrows():
                        rp = r.get("responsible_person","—") or "—"
                        rows.append(
                            f"{r['act_id']} · {r['activity']} · "
                            f"Variance: {r['variance_pct']*100:.1f}% · "
                            f"Slip: {r['weeks_slip']}wks · "
                            f"Reason: {r['delay_reason']} · "
                            f"Responsible: {rp}"
                        )
                    behind_str = "\n".join(rows)

                prompt = f"""You are a senior planning engineer writing formal Minutes of Meeting.

PROJECT: {project_name}
WEEK: {selected_week} | DATE: {meeting_date} | CHAIRED BY: {meeting_chairperson}
BASELINE: {baseline_ver} | CONTRACTOR: {contractor}

DELAYED ACTIVITIES:
{behind_str}

CRITICAL PATH AT RISK:
{critical_acts[['act_id','activity','variance_pct','weeks_slip','delay_reason']].to_string() if not critical_acts.empty else 'None'}

RECURRING ISSUES:
{det_r3['detail'].tolist() if not det_r3.empty else 'None'}

AI SUMMARY: {narrative.get('executive_summary','') if narrative else ''}

Generate formal MOM with sections: {', '.join(sections)}

CRITICAL RULES:
- Every action item MUST include: Action | Responsible Person | Deadline
- Use the responsible person name from the delayed activities list above
- If responsible person is listed, address action directly to them by name
- Deadline format: "Next WPR (Week {selected_week+1})" unless critical
- Recurring issues: name the person responsible and consequence if unresolved
- Be specific — no vague statements
- Use ALL CAPS for section titles

Start with:
MINUTES OF MEETING
Project: {project_name}
Date: {meeting_date}
Week: {selected_week}
Chaired by: {meeting_chairperson}
Attendees: {', '.join([a.strip() for a in attendees.split(chr(10)) if a.strip()])}
"""
                mom_text = call_groq(prompt, max_tokens=2000)

            st.divider()
            st.subheader("Generated MOM")
            st.text_area(
                "", value=mom_text, height=500,
                label_visibility="collapsed")
            st.download_button(
                "⬇️ Download MOM as .txt",
                mom_text,
                f"MOM_SiteInsight_Wk{selected_week:02d}_{meeting_date}.txt",
                "text/plain",
            )

            # ── Email notifications ───────────────────────────
            st.divider()
            st.subheader("📧 Send Email Notifications")
            st.caption(
                "Sends each responsible person only their own action items.")

            if not EMAIL_SENDER:
                st.warning(
                    "EMAIL_SENDER not set in .env file. "
                    "Add it to enable notifications."
                )
            else:
                # Build person → actions map
                person_actions = {}
                if not behind_acts.empty and "responsible_person" in behind_acts.columns:
                    for _, r in behind_acts.iterrows():
                        rp = r.get("responsible_person")
                        if rp and str(rp) not in ["nan","None",""]:
                            person = str(rp)
                            if person not in person_actions:
                                person_actions[person] = []
                            person_actions[person].append(
                                f"• {r['activity']} ({r['act_id']}) — "
                                f"Variance: {r['variance_pct']*100:.1f}%, "
                                f"Slip: {r['weeks_slip']} wks. "
                                f"Reason: {r['delay_reason']}. "
                                f"Deadline: Next WPR (Week {selected_week+1})"
                            )

                if not person_actions:
                    st.info(
                        "No responsible persons assigned yet. "
                        "Add names to Wk11+ Excel files and re-run pipeline."
                    )
                else:
                    st.markdown(
                        f"Found **{len(person_actions)} responsible persons** "
                        f"with action items:"
                    )
                    email_map = {}
                    for person in person_actions:
                        ec1, ec2 = st.columns([1,2])
                        ec1.write(f"👤 **{person}**")
                        email_addr = ec2.text_input(
                            f"Email for {person}",
                            placeholder="name@company.com",
                            key=f"email_{person}",
                            label_visibility="collapsed",
                        )
                        email_map[person] = email_addr

                    if st.button(
                        "📧 Send action items to all responsible persons",
                        type="primary"
                    ):
                        sent_count = 0
                        for person, actions in person_actions.items():
                            recipient = email_map.get(person,"").strip()
                            if not recipient or "@" not in recipient:
                                st.warning(
                                    f"No valid email for {person} — skipped")
                                continue

                            body = f"""Dear {person},

This is your action summary from the Week {selected_week:02d} WPR review meeting for {project_name}.

YOUR ACTION ITEMS:
{''.join([chr(10)+a for a in actions])}

Meeting date: {meeting_date}
Chaired by: {meeting_chairperson}
Deadline: Next WPR — Week {selected_week+1:02d}

Please ensure these items are addressed before the next review.

Regards,
SiteInsight · {project_name}
"""
                            subject = (
                                f"SiteInsight Alert — Your action items "
                                f"from Week {selected_week:02d} review"
                            )
                            if send_email(recipient, subject, body):
                                st.success(
                                    f"✅ Sent to {person} ({recipient})")
                                sent_count += 1

                        if sent_count > 0:
                            st.success(
                                f"Notifications sent to "
                                f"{sent_count} person(s)."
                            )


# ════════════════════════════════════════════════════════════════
# TAB 3 · GENERATE PPT
# ════════════════════════════════════════════════════════════════
with tab3:
    st.subheader("📊 Presentation Generator")
    st.caption(
        "Choose slides to include — SiteInsight builds the deck automatically.")

    with st.form("ppt_form"):
        st.markdown("**Select slides**")
        pc1, pc2 = st.columns(2)
        with pc1:
            slide_cover   = st.checkbox("Cover slide",            value=True)
            slide_health  = st.checkbox("Project health summary", value=True)
            slide_phase   = st.checkbox("Phase-wise status",      value=True)
        with pc2:
            slide_risks   = st.checkbox("Top risks",              value=True)
            slide_actions = st.checkbox("Next week action plan",  value=True)
            slide_all     = st.checkbox("Select ALL",             value=False)

        st.markdown("**Details**")
        ppt_by       = st.text_input("Prepared by", value="Planning Team")
        ppt_audience = st.selectbox(
            "Audience",
            ["Project Director","Client","Internal Team","Site Review Meeting"]
        )
        gen_ppt = st.form_submit_button(
            "Generate Presentation", type="primary")

    if gen_ppt:
        if slide_all:
            slide_cover = slide_health = slide_phase = \
                slide_risks = slide_actions = True

        selected_slides = (
            (["cover"]             if slide_cover   else []) +
            (["project_health"]    if slide_health  else []) +
            (["phase_status"]      if slide_phase   else []) +
            (["top_risks"]         if slide_risks   else []) +
            (["next_week_actions"] if slide_actions else [])
        )
        if not selected_slides:
            st.warning("Please select at least one slide.")
        else:
            with st.spinner("Generating slide content via Groq..."):
                behind_acts = activities[
                    activities["variance_pct"] < -0.05
                ] if not activities.empty else pd.DataFrame()

                phase_summary = ""
                if not activities.empty:
                    for phase, grp in activities.groupby("phase"):
                        avg_v    = grp["variance_pct"].mean() * 100
                        behind_n = (grp["variance_pct"] < -0.05).sum()
                        phase_summary += (
                            f"{phase}: avg variance {avg_v:.1f}%, "
                            f"{behind_n}/{len(grp)} behind\n"
                        )

                prompt = f"""You are preparing a PowerPoint for {ppt_audience} on Week {selected_week} WPR for {project_name}.

PROJECT: {project_name} | WEEK: {selected_week} | BASELINE: {baseline_ver}
Behind plan: {len(behind_acts)}/{len(activities)} activities
Phase summary:
{phase_summary}
Top delayed:
{behind_acts[['act_id','activity','phase','variance_pct','weeks_slip','delay_reason']].head(5).to_string() if not behind_acts.empty else 'None'}
AI SUMMARY: {narrative.get('executive_summary','') if narrative else ''}
KEY RISKS: {narrative.get('key_risks','') if narrative else ''}

Generate for slides: {', '.join(selected_slides)}
Return JSON only — no markdown fences:
{{
  "cover":            {{"title":"...","bullets":[...],"highlight":"..."}},
  "project_health":   {{"title":"...","bullets":[...],"highlight":"..."}},
  "phase_status":     {{"title":"...","bullets":[...],"highlight":"..."}},
  "top_risks":        {{"title":"...","bullets":[...],"highlight":"..."}},
  "next_week_actions":{{"title":"...","bullets":[...],"highlight":"..."}}
}}
Rules: 3-5 bullets max 15 words each, specific, audience is {ppt_audience}"""

                raw = call_groq(prompt, max_tokens=1200)
                try:
                    if "```" in raw:
                        raw = raw.split("```")[1]
                        if raw.startswith("json"):
                            raw = raw[4:]
                    slide_content = json.loads(raw.strip())
                except Exception:
                    slide_content = {}

            st.divider()
            st.subheader("Slide Preview")

            SLIDE_ORDER = [
                "cover","project_health","phase_status",
                "top_risks","next_week_actions"
            ]
            SLIDE_LABELS = {
                "cover":             "Cover",
                "project_health":    "Project Health",
                "phase_status":      "Phase-wise Status",
                "top_risks":         "Top Risks",
                "next_week_actions": "Next Week Actions",
            }

            for key in SLIDE_ORDER:
                if key not in slide_content:
                    continue
                s = slide_content[key]
                with st.expander(
                    f"Slide: {SLIDE_LABELS.get(key,key)}",
                    expanded=True
                ):
                    scol1, scol2 = st.columns([2,1])
                    with scol1:
                        st.markdown(f"**{s.get('title','')}**")
                        for b in s.get("bullets",[]):
                            st.markdown(f"• {b}")
                    with scol2:
                        if s.get("highlight"):
                            st.metric("Key stat", s["highlight"])

            # Build PPTX
            try:
                from pptx import Presentation as PPTXPres
                from pptx.util import Inches, Pt
                from pptx.dml.color import RGBColor as RGB

                prs              = PPTXPres()
                prs.slide_width  = Inches(13.33)
                prs.slide_height = Inches(7.5)
                blank            = prs.slide_layouts[6]

                C_NAVY   = RGB(0x1E,0x27,0x61)
                C_WHITE  = RGB(0xFF,0xFF,0xFF)
                C_ACCENT = RGB(0x00,0xA8,0xE8)
                C_GRAY   = RGB(0x44,0x44,0x44)

                def add_slide(title_text, bullets,
                              highlight=None, is_cover=False):
                    sl  = prs.slides.add_slide(blank)
                    bg  = sl.shapes.add_shape(
                        1,0,0,prs.slide_width,prs.slide_height)
                    bg.fill.solid()
                    bg.fill.fore_color.rgb = C_NAVY if is_cover else C_WHITE
                    bg.line.fill.background()
                    bar = sl.shapes.add_shape(
                        1,0,0,prs.slide_width,Inches(0.07))
                    bar.fill.solid()
                    bar.fill.fore_color.rgb = C_ACCENT
                    bar.line.fill.background()

                    tb = sl.shapes.add_textbox(
                        Inches(0.5),Inches(0.25),Inches(12),Inches(1.2))
                    p  = tb.text_frame.paragraphs[0]
                    p.text          = title_text
                    p.font.size     = Pt(34 if is_cover else 26)
                    p.font.bold     = True
                    p.font.color.rgb = C_WHITE if is_cover else C_NAVY

                    if is_cover:
                        sb = sl.shapes.add_textbox(
                            Inches(0.5),Inches(1.8),Inches(12),Inches(1))
                        sp = sb.text_frame.paragraphs[0]
                        sp.text          = (
                            f"Week {selected_week:02d} · {project_name}")
                        sp.font.size     = Pt(18)
                        sp.font.color.rgb = C_ACCENT
                        db  = sl.shapes.add_textbox(
                            Inches(0.5),Inches(3.0),Inches(12),Inches(2.5))
                        dtf = db.text_frame
                        dtf.word_wrap = True
                        for i, line in enumerate([
                            f"Baseline: {baseline_ver}",
                            f"Contractor: {contractor}",
                            f"Prepared by: {ppt_by}",
                            f"Audience: {ppt_audience}",
                        ]):
                            dp = (dtf.paragraphs[0] if i == 0
                                  else dtf.add_paragraph())
                            dp.text          = line
                            dp.font.size     = Pt(14)
                            dp.font.color.rgb = C_WHITE
                    else:
                        if highlight:
                            hb = sl.shapes.add_shape(
                                1,Inches(10.5),Inches(1.6),
                                Inches(2.3),Inches(1.4))
                            hb.fill.solid()
                            hb.fill.fore_color.rgb = C_ACCENT
                            hb.line.fill.background()
                            ht             = hb.text_frame.paragraphs[0]
                            ht.text        = str(highlight)
                            ht.font.size   = Pt(20)
                            ht.font.bold   = True
                            ht.font.color.rgb = C_WHITE
                        bb  = sl.shapes.add_textbox(
                            Inches(0.5),Inches(1.6),Inches(9.8),Inches(5.5))
                        btf = bb.text_frame
                        btf.word_wrap = True
                        for i, bullet in enumerate(bullets):
                            bp = (btf.paragraphs[0] if i == 0
                                  else btf.add_paragraph())
                            bp.text          = f"▸  {bullet}"
                            bp.font.size     = Pt(15)
                            bp.font.color.rgb = C_GRAY
                            bp.space_after   = Pt(8)

                    ft = sl.shapes.add_textbox(
                        Inches(0.5),Inches(7.1),Inches(12),Inches(0.3))
                    fp = ft.text_frame.paragraphs[0]
                    fp.text = (
                        f"SiteInsight · {project_name} · "
                        f"Week {selected_week:02d} · CONFIDENTIAL"
                    )
                    fp.font.size     = Pt(8)
                    fp.font.color.rgb = RGB(0x99,0x99,0x99)

                for key in SLIDE_ORDER:
                    if key not in slide_content:
                        continue
                    s = slide_content[key]
                    add_slide(
                        s.get("title",key),
                        s.get("bullets",[]),
                        s.get("highlight"),
                        is_cover=(key=="cover"),
                    )

                buf = io.BytesIO()
                prs.save(buf)
                buf.seek(0)
                st.download_button(
                    "⬇️ Download PowerPoint (.pptx)",
                    buf,
                    f"SiteInsight_Wk{selected_week:02d}.pptx",
                    "application/vnd.openxmlformats-officedocument"
                    ".presentationml.presentation",
                )
            except ImportError:
                st.warning(
                    "python-pptx not installed. "
                    "Run: pip install python-pptx"
                )


# ════════════════════════════════════════════════════════════════
# TAB 4 · GENERATE GRAPHS
# ════════════════════════════════════════════════════════════════
with tab4:
    st.subheader("📈 Graph Generator")

    with st.form("graph_form"):
        st.markdown("**Select graphs**")
        gc1, gc2 = st.columns(2)
        with gc1:
            g_variance = st.checkbox(
                "Variance trend — all weeks",       value=True)
            g_phase    = st.checkbox(
                "Phase-wise completion bar chart",  value=True)
            g_slip     = st.checkbox(
                "Weeks slip by activity",           value=True)
        with gc2:
            g_delay    = st.checkbox(
                "Delay reasons breakdown",          value=True)
            g_critical = st.checkbox(
                "Critical path: planned vs actual", value=True)
            g_all      = st.checkbox("Generate ALL graphs", value=False)
        gen_graphs = st.form_submit_button("Generate Graphs", type="primary")

    if gen_graphs:
        if g_all:
            g_variance = g_phase = g_slip = g_delay = g_critical = True

        selected_graphs = (
            (["variance"] if g_variance else []) +
            (["phase"]    if g_phase    else []) +
            (["slip"]     if g_slip     else []) +
            (["delay"]    if g_delay    else []) +
            (["critical"] if g_critical else [])
        )
        if not selected_graphs:
            st.warning("Select at least one graph.")
        else:
            st.divider()

            # 1 · Variance trend
            if "variance" in selected_graphs and not history_df.empty:
                st.markdown("#### Variance trend — all weeks")
                trend = (history_df.groupby("week_number")["variance_pct"]
                         .mean().reset_index())
                trend["variance_pct"] *= 100
                fig, ax = plt.subplots(figsize=(10,3.5))
                colors = [RED if v < 0 else GREEN
                          for v in trend["variance_pct"]]
                ax.bar(trend["week_number"], trend["variance_pct"],
                       color=colors, width=0.6, zorder=3)
                ax.axhline(0, color=GRAY, linewidth=0.8, linestyle="--")
                ax.axvline(selected_week-0.5, color=ACCENT, linewidth=1.5,
                           linestyle="--",
                           label=f"Wk{selected_week} selected")
                ax.set_xlabel("Week"); ax.set_ylabel("Avg Variance (%)")
                ax.set_title(
                    f"{project_name} · Average Variance by Week",
                    fontweight="bold")
                ax.legend(); ax.grid(axis="y", alpha=0.3)
                ax.set_facecolor("#FAFAFA")
                fig.patch.set_facecolor("white"); plt.tight_layout()
                buf = io.BytesIO()
                fig.savefig(buf, format="png", dpi=150, bbox_inches="tight")
                buf.seek(0)
                st.image(buf, use_column_width=True)
                st.download_button(
                    "⬇️ Download", buf,
                    f"variance_Wk{selected_week:02d}.png",
                    "image/png", key="dl_v")
                plt.close(fig); st.divider()

            # 2 · Phase-wise completion
            if "phase" in selected_graphs and not activities.empty:
                st.markdown("#### Phase-wise completion — this week")
                pd_data = activities.groupby("phase").agg(
                    actual=("cum_actual_pct","mean"),
                    planned=("cum_planned_pct","mean")
                ).reset_index()
                pd_data[["actual","planned"]] *= 100
                fig, ax = plt.subplots(figsize=(9,4))
                x = range(len(pd_data)); w = 0.35
                b1 = ax.bar([i-w/2 for i in x], pd_data["planned"],
                            w, label="Planned %", color=ACCENT, alpha=0.85)
                b2 = ax.bar([i+w/2 for i in x], pd_data["actual"],
                            w, label="Actual %", color=NAVY, alpha=0.85)
                ax.set_xticks(list(x))
                ax.set_xticklabels(pd_data["phase"])
                ax.set_ylabel("Completion (%)")
                ax.set_title(
                    f"Phase-wise Status · Week {selected_week:02d}",
                    fontweight="bold")
                ax.legend(); ax.grid(axis="y", alpha=0.3)
                ax.set_facecolor("#FAFAFA")
                fig.patch.set_facecolor("white")
                for b in list(b1)+list(b2):
                    ax.text(
                        b.get_x()+b.get_width()/2,
                        b.get_height()+0.5,
                        f"{b.get_height():.0f}%",
                        ha="center", va="bottom", fontsize=8)
                plt.tight_layout()
                buf = io.BytesIO()
                fig.savefig(buf, format="png", dpi=150, bbox_inches="tight")
                buf.seek(0)
                st.image(buf, use_column_width=True)
                st.download_button(
                    "⬇️ Download", buf,
                    f"phase_Wk{selected_week:02d}.png",
                    "image/png", key="dl_p")
                plt.close(fig); st.divider()

            # 3 · Slip tracker
            if "slip" in selected_graphs and not activities.empty:
                st.markdown("#### Weeks slip by activity")
                slip_data = activities[
                    activities["weeks_slip"] > 0
                ].sort_values("weeks_slip", ascending=True)
                if slip_data.empty:
                    st.success("No activities with weeks slip this week.")
                else:
                    fig, ax = plt.subplots(
                        figsize=(9, max(3.5, len(slip_data)*0.45)))
                    colors = [
                        RED if row.get("is_critical_path") else AMBER
                        for _, row in slip_data.iterrows()
                    ]
                    bars = ax.barh(
                        slip_data["act_id"],
                        slip_data["weeks_slip"],
                        color=colors, height=0.55)
                    ax.set_xlabel("Weeks Slip")
                    ax.set_title(
                        f"Weeks Slip · Week {selected_week:02d}",
                        fontweight="bold")
                    ax.grid(axis="x", alpha=0.3)
                    ax.set_facecolor("#FAFAFA")
                    fig.patch.set_facecolor("white")
                    from matplotlib.patches import Patch
                    ax.legend(handles=[
                        Patch(color=RED,   label="Critical"),
                        Patch(color=AMBER, label="Non-critical")
                    ])
                    for b in bars:
                        ax.text(
                            b.get_width()+0.05,
                            b.get_y()+b.get_height()/2,
                            f"{int(b.get_width())}w",
                            va="center", fontsize=8)
                    plt.tight_layout()
                    buf = io.BytesIO()
                    fig.savefig(
                        buf, format="png", dpi=150, bbox_inches="tight")
                    buf.seek(0)
                    st.image(buf, use_column_width=True)
                    st.download_button(
                        "⬇️ Download", buf,
                        f"slip_Wk{selected_week:02d}.png",
                        "image/png", key="dl_s")
                    plt.close(fig); st.divider()

            # 4 · Delay reasons
            if "delay" in selected_graphs and not activities.empty:
                st.markdown("#### Delay reasons breakdown")
                delayed = activities[
                    (activities["variance_pct"] < 0) &
                    (activities["delay_reason"].notna()) &
                    (activities["delay_reason"] != "No Delay")
                ]
                if delayed.empty:
                    st.success("No delay reasons to display.")
                else:
                    counts  = delayed["delay_reason"].value_counts()
                    fig, ax = plt.subplots(figsize=(7,4.5))
                    palette = [NAVY,ACCENT,RED,AMBER,GREEN,
                               "#7B61FF","#FF6B6B","#4ECDC4"]
                    wedges, texts, autotexts = ax.pie(
                        counts.values,
                        labels=counts.index,
                        autopct="%1.0f%%",
                        colors=palette[:len(counts)],
                        startangle=140,
                        pctdistance=0.82,
                    )
                    for t in texts:
                        t.set_fontsize(9)
                    for at in autotexts:
                        at.set_fontsize(8)
                        at.set_color("white")
                        at.set_fontweight("bold")
                    ax.set_title(
                        f"Delay Reasons · Week {selected_week:02d}",
                        fontweight="bold")
                    fig.patch.set_facecolor("white"); plt.tight_layout()
                    buf = io.BytesIO()
                    fig.savefig(
                        buf, format="png", dpi=150, bbox_inches="tight")
                    buf.seek(0)
                    st.image(buf, use_column_width=True)
                    st.download_button(
                        "⬇️ Download", buf,
                        f"delay_Wk{selected_week:02d}.png",
                        "image/png", key="dl_d")
                    plt.close(fig); st.divider()

            # 5 · Critical path
            if ("critical" in selected_graphs
                    and not activities.empty
                    and "is_critical_path" in activities.columns):
                st.markdown("#### Critical path: planned vs actual")
                cp = activities[
                    activities["is_critical_path"] == True].copy()
                if cp.empty:
                    st.info("No critical path activities found.")
                else:
                    cp = cp.sort_values("variance_pct")
                    fig, ax = plt.subplots(
                        figsize=(9, max(3.5, len(cp)*0.55)))
                    x = range(len(cp))
                    ax.barh(
                        [i+0.2 for i in x],
                        cp["cum_planned_pct"]*100,
                        height=0.35, label="Planned %",
                        color=ACCENT, alpha=0.8)
                    ax.barh(
                        [i-0.2 for i in x],
                        cp["cum_actual_pct"]*100,
                        height=0.35, label="Actual %",
                        color=NAVY, alpha=0.8)
                    ax.set_yticks(list(x))
                    ax.set_yticklabels(cp["act_id"], fontsize=9)
                    ax.set_xlabel("Completion (%)")
                    ax.set_title(
                        f"Critical Path · Week {selected_week:02d}",
                        fontweight="bold")
                    ax.legend(); ax.grid(axis="x", alpha=0.3)
                    ax.set_facecolor("#FAFAFA")
                    fig.patch.set_facecolor("white"); plt.tight_layout()
                    buf = io.BytesIO()
                    fig.savefig(
                        buf, format="png", dpi=150, bbox_inches="tight")
                    buf.seek(0)
                    st.image(buf, use_column_width=True)
                    st.download_button(
                        "⬇️ Download", buf,
                        f"critical_Wk{selected_week:02d}.png",
                        "image/png", key="dl_c")
                    plt.close(fig)

            st.success(
                f"Generated {len(selected_graphs)} graph(s) successfully.")


# ════════════════════════════════════════════════════════════════
# TAB 5 · RECOVERY SIMULATOR
# ════════════════════════════════════════════════════════════════
with tab5:
    st.subheader("🔄 Recovery Simulator")
    st.caption(
        "Select a delayed activity and simulate catch-up scenarios "
        "based on different acceleration targets."
    )

    if activities.empty or history_df.empty:
        st.info("No activity data available. Load history first.")
    else:
        behind_acts = activities[activities["variance_pct"] < -0.05].copy()
        if behind_acts.empty:
            st.success("No delayed activities this week.")
        else:
            act_options = [
                f"{row['act_id']} · {row['activity']} "
                f"({row['variance_pct']*100:.1f}%)"
                for _, row in behind_acts.iterrows()
            ]
            selected_act_label = st.selectbox(
                "Select delayed activity to simulate", act_options)
            selected_act_id    = selected_act_label.split(" · ")[0]
            act_row = behind_acts[
                behind_acts["act_id"] == selected_act_id].iloc[0]

            current_pct   = float(act_row["cum_actual_pct"])
            planned_pct   = float(act_row["cum_planned_pct"])
            variance      = float(act_row["variance_pct"])
            weeks_slip    = int(act_row["weeks_slip"])
            delay_reason  = act_row.get("delay_reason","—")
            resp_person   = act_row.get("responsible_person","—") or "—"
            hist_velocity = compute_weekly_velocity(
                history_df, selected_act_id)

            st.divider()
            col_s1, col_s2, col_s3, col_s4, col_s5 = st.columns(5)
            col_s1.metric("Current actual", f"{current_pct*100:.1f}%")
            col_s2.metric("Planned target", f"{planned_pct*100:.1f}%")
            col_s3.metric("Variance",       f"{variance*100:.1f}%",
                          delta_color="inverse")
            col_s4.metric("Weeks slip",     f"{weeks_slip} wks",
                          delta_color="inverse")
            col_s5.metric("Responsible",    resp_person)
            st.caption(
                f"Delay reason: **{delay_reason}** · "
                f"Historical rate: **{hist_velocity*100:.2f}%/wk**"
            )

            st.divider()
            sc1, sc2 = st.columns(2)
            with sc1:
                acceleration = st.slider(
                    "Target weekly progress rate (%)",
                    min_value=1.0, max_value=25.0,
                    value=max(float(hist_velocity*100*1.5), 5.0),
                    step=0.5,
                )
            with sc2:
                total_project_weeks = st.number_input(
                    "Total project weeks",
                    min_value=20, max_value=200, value=80)

            accel_rate        = acceleration / 100
            realistic_rate    = hist_velocity if hist_velocity > 0 else 0.02
            optimistic_rate   = accel_rate
            conservative_rate = accel_rate * 0.6

            proj_realistic    = compute_projected_finish(
                current_pct, selected_week,
                realistic_rate, total_project_weeks)
            proj_optimistic   = compute_projected_finish(
                current_pct, selected_week,
                optimistic_rate, total_project_weeks)
            proj_conservative = compute_projected_finish(
                current_pct, selected_week,
                conservative_rate, total_project_weeks)

            st.divider()
            st.markdown("**Projected finish week**")
            r1, r2, r3 = st.columns(3)
            r1.metric("Conservative",
                      f"Week {min(proj_conservative, total_project_weeks)}",
                      delta_color="inverse")
            r2.metric("Target (your rate)",
                      f"Week {min(proj_optimistic, total_project_weeks)}",
                      delta_color="normal")
            r3.metric("At current rate",
                      f"Week {min(proj_realistic, total_project_weeks)}",
                      delta_color="inverse")

            st.divider()
            future_weeks = list(range(
                selected_week,
                min(selected_week+30, total_project_weeks+1)))

            def project_curve(start_pct, rate, weeks):
                pcts = []; p = start_pct
                for _ in weeks:
                    p = min(p+rate, 1.0); pcts.append(p*100)
                return pcts

            fig, ax = plt.subplots(figsize=(10,4))
            ax.plot(future_weeks,
                    project_curve(current_pct, optimistic_rate, future_weeks),
                    color=GREEN, linewidth=2,
                    label=f"Target ({acceleration:.1f}%/wk)")
            ax.plot(future_weeks,
                    project_curve(current_pct, conservative_rate, future_weeks),
                    color=AMBER, linewidth=1.5, linestyle="--",
                    label=f"Conservative ({acceleration*0.6:.1f}%/wk)")
            ax.plot(future_weeks,
                    project_curve(current_pct, realistic_rate, future_weeks),
                    color=RED, linewidth=1.5, linestyle=":",
                    label=f"Current rate ({hist_velocity*100:.1f}%/wk)")
            ax.axhline(100, color=GRAY, linewidth=0.8,
                       linestyle="--", label="100% complete")
            ax.axhline(planned_pct*100, color=NAVY, linewidth=0.8,
                       linestyle="-.",
                       label=f"Planned ({planned_pct*100:.0f}%)")
            ax.set_xlabel("Week"); ax.set_ylabel("Cumulative Progress (%)")
            ax.set_title(
                f"Recovery Simulation · {selected_act_id} · "
                f"{act_row['activity']}",
                fontweight="bold")
            ax.legend(fontsize=9); ax.grid(alpha=0.3)
            ax.set_facecolor("#FAFAFA")
            fig.patch.set_facecolor("white"); plt.tight_layout()
            buf = io.BytesIO()
            fig.savefig(buf, format="png", dpi=150, bbox_inches="tight")
            buf.seek(0)
            st.image(buf, use_column_width=True)
            st.download_button(
                "⬇️ Download simulation chart", buf,
                f"recovery_{selected_act_id}_Wk{selected_week:02d}.png",
                "image/png", key="dl_sim")
            plt.close(fig)

            st.divider()
            if st.button(
                "Generate AI recovery recommendation", type="primary"
            ):
                with st.spinner("Analysing via Groq..."):
                    prompt = f"""You are a senior construction planning consultant.
ACTIVITY: {act_row['activity']} ({selected_act_id})
PHASE: {act_row['phase']} | RESPONSIBLE: {resp_person}
CURRENT: {current_pct*100:.1f}% | PLANNED: {planned_pct*100:.1f}%
VARIANCE: {variance*100:.1f}% | WEEKS SLIP: {weeks_slip}
DELAY REASON: {delay_reason}
HISTORICAL RATE: {hist_velocity*100:.2f}%/wk | TARGET: {acceleration:.1f}%/wk
At current rate: Week {proj_realistic} | At target: Week {proj_optimistic}

Write a specific recovery recommendation (max 200 words):
1. Is {acceleration:.1f}%/week realistic given the delay reason?
2. What must {resp_person} do THIS WEEK specifically?
3. Consequence if not recovered by Week {selected_week+1}?
4. One direct instruction to the contractor.
Direct construction language only."""
                    rec = call_groq(prompt, max_tokens=400)
                st.markdown("**AI Recovery Recommendation**")
                st.warning(rec)


# ════════════════════════════════════════════════════════════════
# TAB 6 · EARLY WARNING
# ════════════════════════════════════════════════════════════════
with tab6:
    st.subheader("⚠️ Early Warning System")
    st.caption(
        "Identifies activities likely to slip NEXT week based on "
        "trend velocity — before it appears in the WPR."
    )

    LOOKBACK = 4

    if history_df.empty or activities.empty:
        st.info("Need at least 3 weeks of history for trend analysis.")
    else:
        warnings_list = []
        for _, act in activities.iterrows():
            aid      = act["act_id"]
            variance = float(act.get("variance_pct", 0) or 0)
            critical = bool(act.get("is_critical_path", False))
            slip     = int(act.get("weeks_slip", 0) or 0)
            resp     = act.get("responsible_person","—") or "—"
            slope    = compute_trend_slope(history_df, aid, LOOKBACK)

            risk_score = 0; risk_flags = []
            if slope < -0.015:
                risk_score += 3
                risk_flags.append(
                    f"Variance worsening "
                    f"{abs(slope)*100:.1f}%/wk over {LOOKBACK} wks")
            if variance < -0.05:
                risk_score += 2
                risk_flags.append(
                    f"Already {abs(variance)*100:.1f}% behind plan")
            if slip >= 2:
                risk_score += 2
                risk_flags.append(f"{slip} weeks slip accumulated")
            if critical:
                risk_score += 2
                risk_flags.append("On critical path")
            if slope < 0 and variance < -0.03:
                risk_score += 1
                risk_flags.append("Negative trend + negative variance")

            if risk_score >= 3:
                risk_level = (
                    "HIGH"   if risk_score >= 7 else
                    "MEDIUM" if risk_score >= 4 else "LOW"
                )
                warnings_list.append({
                    "act_id":     aid,
                    "activity":   act["activity"],
                    "phase":      act["phase"],
                    "variance":   variance,
                    "slope":      slope,
                    "slip":       slip,
                    "critical":   critical,
                    "responsible": resp,
                    "risk_score": risk_score,
                    "risk_level": risk_level,
                    "flags":      risk_flags,
                })

        warnings_list.sort(key=lambda x: x["risk_score"], reverse=True)

        if not warnings_list:
            st.success("No early warning signals detected.")
        else:
            high   = [w for w in warnings_list if w["risk_level"]=="HIGH"]
            medium = [w for w in warnings_list if w["risk_level"]=="MEDIUM"]
            low    = [w for w in warnings_list if w["risk_level"]=="LOW"]

            w1, w2, w3 = st.columns(3)
            w1.metric("HIGH risk", len(high),
                      delta="Immediate action" if high else None,
                      delta_color="inverse")
            w2.metric("MEDIUM risk", len(medium),
                      delta="Monitor closely" if medium else None,
                      delta_color="inverse")
            w3.metric("LOW risk", len(low))
            st.divider()

            level_icons = {"HIGH":"🔴","MEDIUM":"🟡","LOW":"🟢"}
            for w in warnings_list:
                with st.expander(
                    f"{level_icons[w['risk_level']]} {w['risk_level']} · "
                    f"{w['act_id']} · {w['activity']} "
                    f"(Score: {w['risk_score']})",
                    expanded=(w["risk_level"]=="HIGH"),
                ):
                    e1, e2, e3, e4, e5 = st.columns(5)
                    e1.metric("Variance",
                              f"{w['variance']*100:.1f}%",
                              delta_color="inverse")
                    e2.metric("Trend slope",
                              f"{w['slope']*100:.2f}%/wk",
                              delta_color=(
                                  "inverse" if w["slope"] < 0
                                  else "normal"))
                    e3.metric("Weeks slip",
                              f"{w['slip']} wks",
                              delta_color=(
                                  "inverse" if w["slip"] > 0
                                  else "normal"))
                    e4.metric("Critical",
                              "YES" if w["critical"] else "No")
                    e5.metric("Responsible", w["responsible"])

                    st.markdown("**Why flagged:**")
                    for flag in w["flags"]:
                        st.markdown(f"• {flag}")

            # Trend chart
            st.divider()
            st.markdown("#### Variance trend — at-risk activities")
            at_risk_ids = [w["act_id"] for w in warnings_list[:6]]
            hist_subset = history_df[
                history_df["act_id"].isin(at_risk_ids)]

            if not hist_subset.empty:
                fig, ax = plt.subplots(figsize=(10,4))
                for aid in at_risk_ids:
                    act_data   = hist_subset[
                        hist_subset["act_id"]==aid
                    ].sort_values("week_number")
                    if len(act_data) < 2:
                        continue
                    risk_level = next(
                        (w["risk_level"] for w in warnings_list
                         if w["act_id"]==aid), "LOW")
                    color = (RED   if risk_level=="HIGH" else
                             AMBER if risk_level=="MEDIUM" else GRAY)
                    ax.plot(
                        act_data["week_number"],
                        act_data["variance_pct"]*100,
                        label=aid,
                        linewidth=2.5 if risk_level=="HIGH" else 1.5,
                        color=color, marker="o", markersize=3)
                ax.axhline(0,   color=GRAY, linewidth=0.8, linestyle="--")
                ax.axhline(-5,  color=RED,  linewidth=0.5,
                           linestyle=":", alpha=0.5)
                ax.axhline(-10, color=RED,  linewidth=0.8,
                           linestyle=":", alpha=0.7)
                ax.axvline(selected_week, color=ACCENT,
                           linewidth=1, linestyle="--", alpha=0.7)
                ax.set_xlabel("Week"); ax.set_ylabel("Variance (%)")
                ax.set_title(
                    "Variance Trajectory — At-risk Activities",
                    fontweight="bold")
                ax.legend(fontsize=8, ncol=3)
                ax.grid(alpha=0.3); ax.set_facecolor("#FAFAFA")
                fig.patch.set_facecolor("white"); plt.tight_layout()
                buf = io.BytesIO()
                fig.savefig(
                    buf, format="png", dpi=150, bbox_inches="tight")
                buf.seek(0)
                st.image(buf, use_column_width=True)
                st.download_button(
                    "⬇️ Download early warning chart", buf,
                    f"early_warning_Wk{selected_week:02d}.png",
                    "image/png", key="dl_ew")
                plt.close(fig)

            # Proactive brief
            st.divider()
            if st.button(
                "Generate proactive weekly brief", type="primary"
            ):
                with st.spinner("Generating via Groq..."):
                    high_names = [
                        f"{w['act_id']} ({w['activity']}, "
                        f"responsible: {w['responsible']}, "
                        f"slope {w['slope']*100:.2f}%/wk)"
                        for w in high
                    ]
                    med_names = [
                        f"{w['act_id']} ({w['activity']}, "
                        f"responsible: {w['responsible']})"
                        for w in medium
                    ]
                    prompt = f"""You are a proactive construction project intelligence system.
Week {selected_week} of {project_name}.

EARLY WARNING ({LOOKBACK}-week trend analysis):
HIGH RISK — will worsen next week: {high_names if high_names else 'None'}
MEDIUM RISK — monitor: {med_names if med_names else 'None'}

Write a PROACTIVE brief (max 250 words):
1. Predict what happens in Week {selected_week+1} if no action taken
2. Name each responsible person and their specific required action THIS WEEK
3. Identify milestone / critical path impact
4. End with ONE thing the PM must do before next site meeting

This is intelligence — tell the team what WILL happen, not what HAS happened.
Direct, action-oriented construction language."""
                    brief = call_groq(prompt, max_tokens=500)
                st.markdown("**Proactive Weekly Brief — Week ahead**")
                st.error(brief)
